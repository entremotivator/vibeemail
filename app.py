import streamlit as st
import json
import re
import time
import io
import base64
import zipfile
import os
import threading
from datetime import datetime
from pathlib import Path

# ─────────────────────────────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="PageForge PRO — AI Document Studio",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ─────────────────────────────────────────────────────────────────
# GLOBAL CSS — Dark brutalist studio aesthetic
# ─────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Bebas+Neue&family=Syne:wght@400;500;600;700;800&family=Space+Mono:ital,wght@0,400;0,700;1,400&family=DM+Sans:wght@300;400;500;600&display=swap');

:root {
    --bg: #060810;
    --surface: #0c0e1a;
    --surface2: #111320;
    --surface3: #161929;
    --border: #1a1e32;
    --border2: #242840;
    --red: #ff3355;
    --orange: #ff7040;
    --blue: #4d6fff;
    --teal: #00c8a0;
    --gold: #ffb830;
    --purple: #8855ff;
    --text: #e4e7f0;
    --muted: #525870;
    --muted2: #3a3f58;
    --success: #2dd4a0;
    --warning: #f59e0b;
    --error: #ff4466;
    --agent-content: #ff3355;
    --agent-pptx: #4d6fff;
    --agent-pdf: #ff7040;
    --agent-seo: #00c8a0;
    --agent-design: #8855ff;
    --agent-qa: #ffb830;
    --agent-copy: #2dd4a0;
}

html, body, [class*="css"] {
    font-family: 'DM Sans', sans-serif !important;
    background: var(--bg) !important;
    color: var(--text) !important;
}
.stApp { background: var(--bg) !important; }
#MainMenu, footer, header { visibility: hidden; }

/* ── SIDEBAR ── */
section[data-testid="stSidebar"] {
    background: var(--surface) !important;
    border-right: 1px solid var(--border) !important;
    width: 310px !important;
}
section[data-testid="stSidebar"] > div { padding: 0 !important; }

/* ── TYPOGRAPHY ── */
h1, h2, h3 {
    font-family: 'Bebas Neue', sans-serif !important;
    letter-spacing: 0.05em !important;
    color: var(--text) !important;
}

/* ── INPUTS ── */
.stTextInput > div > div > input,
.stTextArea > div > div > textarea,
.stNumberInput > div > div > input {
    background: var(--surface2) !important;
    border: 1px solid var(--border2) !important;
    border-radius: 6px !important;
    color: var(--text) !important;
    font-family: 'DM Sans', sans-serif !important;
    font-size: 0.875rem !important;
    caret-color: var(--red) !important;
}
.stTextInput > div > div > input:focus,
.stTextArea > div > div > textarea:focus {
    border-color: var(--red) !important;
    box-shadow: 0 0 0 2px rgba(255,51,85,0.12) !important;
    outline: none !important;
}
label, .stTextInput label, .stTextArea label, .stSelectbox label,
.stNumberInput label, .stSlider label, .stCheckbox label {
    color: var(--muted) !important;
    font-size: 0.7rem !important;
    font-weight: 600 !important;
    letter-spacing: 0.14em !important;
    text-transform: uppercase !important;
    font-family: 'Space Mono', monospace !important;
}

/* ── PRIMARY BUTTONS ── */
.stButton > button {
    background: var(--red) !important;
    color: #fff !important;
    font-family: 'Syne', sans-serif !important;
    font-weight: 700 !important;
    font-size: 0.875rem !important;
    letter-spacing: 0.08em !important;
    text-transform: uppercase !important;
    padding: 0.7rem 1.5rem !important;
    border: none !important;
    border-radius: 6px !important;
    cursor: pointer !important;
    transition: all 0.15s ease !important;
    width: 100%;
}
.stButton > button:hover {
    background: #ff5577 !important;
    transform: translateY(-1px) !important;
    box-shadow: 0 4px 16px rgba(255,51,85,0.4) !important;
}
.stButton > button:disabled {
    background: var(--muted2) !important;
    color: var(--muted) !important;
    transform: none !important;
    box-shadow: none !important;
}

/* ── DOWNLOAD BUTTONS ── */
.stDownloadButton > button {
    background: var(--surface2) !important;
    color: var(--text) !important;
    border: 1px solid var(--border2) !important;
    border-radius: 6px !important;
    font-family: 'Syne', sans-serif !important;
    font-weight: 600 !important;
    font-size: 0.8rem !important;
    text-transform: uppercase !important;
    letter-spacing: 0.06em !important;
    padding: 0.6rem 1rem !important;
    width: 100%;
    transition: all 0.15s ease !important;
}
.stDownloadButton > button:hover {
    border-color: var(--red) !important;
    color: var(--red) !important;
    background: rgba(255,51,85,0.05) !important;
}

/* ── TABS ── */
.stTabs [data-baseweb="tab-list"] {
    background: var(--surface) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    padding: 4px !important;
    gap: 2px !important;
}
.stTabs [data-baseweb="tab"] {
    color: var(--muted) !important;
    border-radius: 6px !important;
    font-family: 'Space Mono', monospace !important;
    font-size: 0.72rem !important;
    padding: 0.5rem 0.9rem !important;
    letter-spacing: 0.06em !important;
}
.stTabs [aria-selected="true"] {
    background: var(--red) !important;
    color: #fff !important;
    font-weight: 700 !important;
}

/* ── SELECT / RADIO / CHECKBOX ── */
.stSelectbox [data-baseweb="select"] > div {
    background: var(--surface2) !important;
    border: 1px solid var(--border2) !important;
    border-radius: 6px !important;
    color: var(--text) !important;
}
.stRadio > div { gap: 0.4rem; }
.stRadio [data-testid="stMarkdownContainer"] p {
    color: var(--text) !important;
    font-family: 'DM Sans', sans-serif !important;
    font-size: 0.875rem !important;
}

/* ── ALERTS ── */
.stSuccess > div {
    background: rgba(45,212,160,0.1) !important;
    border: 1px solid rgba(45,212,160,0.3) !important;
    border-radius: 8px !important;
    color: var(--success) !important;
}
.stError > div {
    background: rgba(255,51,85,0.1) !important;
    border: 1px solid rgba(255,51,85,0.3) !important;
    border-radius: 8px !important;
    color: #ff6680 !important;
}
.stInfo > div {
    background: rgba(77,111,255,0.1) !important;
    border: 1px solid rgba(77,111,255,0.3) !important;
    border-radius: 8px !important;
    color: #7a95ff !important;
}
.stWarning > div {
    background: rgba(245,158,11,0.1) !important;
    border: 1px solid rgba(245,158,11,0.3) !important;
    border-radius: 8px !important;
    color: #fbbf44 !important;
}

/* ── EXPANDER ── */
.streamlit-expanderHeader {
    background: var(--surface2) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    color: var(--text) !important;
    font-family: 'Syne', sans-serif !important;
    font-weight: 600 !important;
    font-size: 0.875rem !important;
}
.streamlit-expanderContent {
    background: var(--surface2) !important;
    border: 1px solid var(--border) !important;
    border-top: none !important;
    border-radius: 0 0 8px 8px !important;
}

/* ── PROGRESS ── */
.stProgress > div > div { background: var(--red) !important; }
.stProgress > div { background: var(--border2) !important; }

/* ── CODE ── */
.stCodeBlock {
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
}
code {
    background: var(--surface2) !important;
    border: 1px solid var(--border) !important;
    border-radius: 4px !important;
    color: var(--gold) !important;
    font-family: 'Space Mono', monospace !important;
    font-size: 0.8rem !important;
    padding: 1px 5px !important;
}

/* ── DIVIDER ── */
hr { border-color: var(--border) !important; margin: 1rem 0 !important; }

/* ── METRIC ── */
[data-testid="metric-container"] {
    background: var(--surface2) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    padding: 0.75rem !important;
}
[data-testid="metric-container"] label {
    color: var(--muted) !important;
    font-family: 'Space Mono', monospace !important;
    font-size: 0.65rem !important;
    letter-spacing: 0.12em !important;
}
[data-testid="metric-container"] [data-testid="stMetricValue"] {
    color: var(--text) !important;
    font-family: 'Bebas Neue', sans-serif !important;
    font-size: 1.6rem !important;
}

/* ── CUSTOM HTML ELEMENTS ── */
.pf-header {
    padding: 1.5rem 1.25rem 1rem;
    border-bottom: 1px solid var(--border);
    background: linear-gradient(180deg, rgba(255,51,85,0.05) 0%, transparent 100%);
}
.pf-logo {
    font-family: 'Bebas Neue', sans-serif;
    font-size: 1.6rem;
    letter-spacing: 0.08em;
    color: var(--text);
    display: flex;
    align-items: center;
    gap: 6px;
    line-height: 1;
}
.pf-logo span { color: var(--red); }
.pf-logo sub {
    font-family: 'Space Mono', monospace;
    font-size: 0.55rem;
    color: var(--muted);
    letter-spacing: 0.15em;
    text-transform: uppercase;
    vertical-align: bottom;
    margin-left: 4px;
    margin-bottom: 2px;
}
.pf-tagline {
    font-family: 'Space Mono', monospace;
    font-size: 0.6rem;
    color: var(--muted2);
    letter-spacing: 0.18em;
    text-transform: uppercase;
    margin-top: 4px;
}

.chip {
    display: inline-flex;
    align-items: center;
    gap: 5px;
    border-radius: 4px;
    padding: 0.18rem 0.6rem;
    font-family: 'Space Mono', monospace;
    font-size: 0.62rem;
    letter-spacing: 0.1em;
    text-transform: uppercase;
    font-weight: 700;
    margin-bottom: 0.75rem;
}
.chip-red { background: rgba(255,51,85,0.12); border: 1px solid rgba(255,51,85,0.3); color: var(--red); }
.chip-blue { background: rgba(77,111,255,0.12); border: 1px solid rgba(77,111,255,0.3); color: var(--blue); }
.chip-orange { background: rgba(255,112,64,0.12); border: 1px solid rgba(255,112,64,0.3); color: var(--orange); }
.chip-teal { background: rgba(0,200,160,0.12); border: 1px solid rgba(0,200,160,0.3); color: var(--teal); }
.chip-gold { background: rgba(255,184,48,0.12); border: 1px solid rgba(255,184,48,0.3); color: var(--gold); }
.chip-purple { background: rgba(136,85,255,0.12); border: 1px solid rgba(136,85,255,0.3); color: var(--purple); }
.chip-success { background: rgba(45,212,160,0.12); border: 1px solid rgba(45,212,160,0.3); color: var(--success); }

.section-label {
    font-family: 'Space Mono', monospace;
    font-size: 0.62rem;
    letter-spacing: 0.16em;
    text-transform: uppercase;
    color: var(--muted);
    margin-bottom: 0.6rem;
    display: block;
}

.page-title {
    font-family: 'Bebas Neue', sans-serif;
    font-size: 2.8rem;
    letter-spacing: 0.05em;
    color: var(--text);
    line-height: 1;
    margin-bottom: 4px;
}
.page-title .accent { color: var(--red); }

.stat-grid-3 {
    display: grid;
    grid-template-columns: 1fr 1fr 1fr;
    gap: 8px;
    margin: 0.75rem 0;
}
.stat-box {
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 6px;
    padding: 0.65rem 0.5rem;
    text-align: center;
}
.stat-box .n {
    font-family: 'Bebas Neue', sans-serif;
    font-size: 1.5rem;
    letter-spacing: 0.04em;
    line-height: 1;
    color: var(--red);
}
.stat-box .l {
    font-family: 'Space Mono', monospace;
    font-size: 0.55rem;
    color: var(--muted);
    text-transform: uppercase;
    letter-spacing: 0.1em;
    margin-top: 3px;
}

.tip-block {
    background: rgba(77,111,255,0.06);
    border-left: 3px solid var(--blue);
    border-radius: 0 6px 6px 0;
    padding: 0.7rem 1rem;
    margin: 0.6rem 0;
    font-size: 0.8rem;
    color: #7a95ff;
    line-height: 1.6;
    font-family: 'DM Sans', sans-serif;
}

.agent-card {
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 0.75rem 0.9rem;
    margin-bottom: 6px;
    transition: border-color 0.3s;
}
.agent-card.running { border-color: var(--warning); box-shadow: 0 0 12px rgba(245,158,11,0.15); }
.agent-card.done { border-color: var(--success); }
.agent-card.error { border-color: var(--error); }
.agent-name {
    font-family: 'Syne', sans-serif;
    font-weight: 700;
    font-size: 0.8rem;
    display: flex;
    align-items: center;
    gap: 6px;
    justify-content: space-between;
}
.agent-role {
    font-size: 0.72rem;
    color: var(--muted);
    margin-top: 2px;
    font-family: 'DM Sans', sans-serif;
}
.agent-tools {
    display: flex;
    flex-wrap: wrap;
    gap: 3px;
    margin-top: 5px;
}
.tool-tag {
    background: var(--surface3);
    border: 1px solid var(--border);
    border-radius: 3px;
    padding: 1px 6px;
    font-family: 'Space Mono', monospace;
    font-size: 0.55rem;
    color: var(--muted);
    letter-spacing: 0.05em;
}

.status-dot {
    width: 7px;
    height: 7px;
    border-radius: 50%;
    display: inline-block;
    flex-shrink: 0;
}
.status-idle { background: var(--muted2); }
.status-running { background: var(--warning); box-shadow: 0 0 6px var(--warning); animation: pulse 1s ease-in-out infinite; }
.status-done { background: var(--success); }
.status-error { background: var(--error); }

@keyframes pulse { 0%,100%{opacity:0.5} 50%{opacity:1} }

.log-entry {
    display: flex;
    gap: 10px;
    align-items: flex-start;
    padding: 5px 8px;
    border-radius: 5px;
    margin-bottom: 2px;
    font-family: 'Space Mono', monospace;
    font-size: 0.7rem;
    line-height: 1.5;
}
.log-entry.running { background: rgba(245,158,11,0.05); }
.log-entry.done { background: rgba(45,212,160,0.04); }
.log-entry.error { background: rgba(255,51,85,0.06); }
.log-entry.info { background: transparent; }
.log-ts { color: var(--muted2); flex-shrink: 0; font-size: 0.62rem; margin-top: 1px; }
.log-agent { font-weight: 700; flex-shrink: 0; font-size: 0.7rem; }
.log-msg { color: var(--text); flex: 1; }

.export-card {
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 10px;
    padding: 1.1rem;
    height: 100%;
    display: flex;
    flex-direction: column;
    gap: 6px;
}
.export-card-icon { font-size: 1.6rem; margin-bottom: 2px; }
.export-card-title {
    font-family: 'Bebas Neue', sans-serif;
    font-size: 1.1rem;
    letter-spacing: 0.06em;
    color: var(--text);
}
.export-card-desc {
    font-size: 0.78rem;
    color: var(--muted);
    line-height: 1.55;
    flex: 1;
}
.export-card-meta {
    font-family: 'Space Mono', monospace;
    font-size: 0.62rem;
    color: var(--muted2);
    margin-bottom: 6px;
}

.slide-preview {
    background: #1a1e2e;
    border: 1px solid var(--border);
    border-radius: 6px;
    padding: 10px;
    aspect-ratio: 16/9;
    display: flex;
    flex-direction: column;
    justify-content: center;
    overflow: hidden;
    position: relative;
}
.slide-preview-num {
    position: absolute;
    top: 6px;
    right: 8px;
    font-family: 'Space Mono', monospace;
    font-size: 0.55rem;
    color: var(--muted2);
}
.slide-preview-title {
    font-family: 'Bebas Neue', sans-serif;
    font-size: 0.9rem;
    letter-spacing: 0.04em;
    color: var(--text);
    margin-bottom: 5px;
    line-height: 1.1;
}
.slide-preview-body {
    font-size: 0.6rem;
    color: var(--muted);
    line-height: 1.4;
}
.slide-preview-bar {
    position: absolute;
    bottom: 0; left: 0; right: 0;
    height: 3px;
}

.pdf-doc-card {
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 1rem;
    transition: border-color 0.2s;
}
.pdf-doc-card:hover { border-color: var(--border2); }
.pdf-doc-type {
    font-family: 'Space Mono', monospace;
    font-size: 0.6rem;
    letter-spacing: 0.14em;
    text-transform: uppercase;
    color: var(--orange);
    margin-bottom: 5px;
}
.pdf-doc-title {
    font-family: 'Syne', sans-serif;
    font-weight: 700;
    font-size: 0.875rem;
    color: var(--text);
    margin-bottom: 4px;
}
.pdf-doc-subtitle {
    font-size: 0.78rem;
    color: var(--muted);
    line-height: 1.5;
}

.feature-edit-card {
    background: var(--surface3);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 0.75rem;
    margin-bottom: 8px;
}

.preview-browser-bar {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 8px 8px 0 0;
    padding: 8px 14px;
    display: flex;
    align-items: center;
    gap: 8px;
    border-bottom: none;
}
.browser-dot { width: 10px; height: 10px; border-radius: 50%; }
.browser-url {
    flex: 1;
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 5px;
    padding: 3px 10px;
    font-family: 'Space Mono', monospace;
    font-size: 0.65rem;
    color: var(--muted);
    margin-left: 8px;
}

.step-block {
    display: flex;
    gap: 12px;
    align-items: flex-start;
    margin-bottom: 10px;
    padding: 10px 0;
    border-bottom: 1px solid var(--border);
}
.step-num {
    font-family: 'Bebas Neue', sans-serif;
    font-size: 2.2rem;
    color: rgba(255,51,85,0.2);
    line-height: 1;
    min-width: 2rem;
    text-align: center;
}
.step-content { flex: 1; padding-top: 2px; }
.step-title {
    font-family: 'Syne', sans-serif;
    font-weight: 700;
    font-size: 0.875rem;
    color: var(--text);
    margin-bottom: 3px;
}
.step-desc { font-size: 0.78rem; color: var(--muted); line-height: 1.5; }

.gen-progress-bar {
    height: 3px;
    background: var(--border2);
    border-radius: 3px;
    margin: 8px 0;
    overflow: hidden;
}
.gen-progress-fill {
    height: 100%;
    border-radius: 3px;
    background: linear-gradient(90deg, var(--red), var(--orange));
    transition: width 0.4s ease;
}

.docs-summary-row {
    display: flex;
    gap: 8px;
    margin: 0.75rem 0;
}
.doc-badge {
    flex: 1;
    background: var(--surface2);
    border: 1px solid var(--border);
    border-radius: 6px;
    padding: 8px 6px;
    text-align: center;
}
.doc-badge-icon { font-size: 1.2rem; margin-bottom: 4px; }
.doc-badge-label {
    font-family: 'Space Mono', monospace;
    font-size: 0.58rem;
    color: var(--muted);
    text-transform: uppercase;
    letter-spacing: 0.08em;
}
.doc-badge-count {
    font-family: 'Bebas Neue', sans-serif;
    font-size: 1.2rem;
    color: var(--text);
    line-height: 1;
    margin-bottom: 2px;
}

.html-code-peek {
    background: var(--surface3);
    border: 1px solid var(--border);
    border-radius: 6px;
    padding: 10px 12px;
    font-family: 'Space Mono', monospace;
    font-size: 0.7rem;
    color: var(--muted);
    max-height: 150px;
    overflow: hidden;
    position: relative;
}
.html-code-peek::after {
    content: '';
    position: absolute;
    bottom: 0; left: 0; right: 0;
    height: 50px;
    background: linear-gradient(transparent, var(--surface3));
}

.theme-swatch {
    width: 100%;
    height: 32px;
    border-radius: 5px;
    cursor: pointer;
    border: 2px solid transparent;
    transition: border-color 0.15s;
    display: flex;
    overflow: hidden;
}
.theme-swatch:hover { border-color: var(--red); }
.theme-swatch div { flex: 1; }

.changelog-item {
    display: flex;
    gap: 10px;
    padding: 8px 0;
    border-bottom: 1px solid var(--border);
    font-size: 0.78rem;
    color: var(--muted);
}
.changelog-tag {
    font-family: 'Space Mono', monospace;
    font-size: 0.6rem;
    padding: 1px 6px;
    border-radius: 3px;
    flex-shrink: 0;
    margin-top: 1px;
}

@keyframes slideIn { from { opacity:0; transform:translateY(8px); } to { opacity:1; transform:translateY(0); } }
.slide-in { animation: slideIn 0.3s ease forwards; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────
# ════════════════════════════════════════════════════════════════
# BMADD AGENT FRAMEWORK
# ════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────

class BmAddTool:
    """Represents a single callable tool within an agent"""
    def __init__(self, name: str, description: str, fn=None):
        self.name = name
        self.description = description
        self.fn = fn
        self.calls = 0
        self.last_result = None
    
    def invoke(self, *args, **kwargs):
        self.calls += 1
        if self.fn:
            self.last_result = self.fn(*args, **kwargs)
            return self.last_result
        return None


class BmAddMemory:
    """Short-term working memory for an agent"""
    def __init__(self, capacity: int = 20):
        self.capacity = capacity
        self._store = []
    
    def write(self, key: str, value):
        self._store.append({"key": key, "value": value, "ts": time.time()})
        if len(self._store) > self.capacity:
            self._store.pop(0)
    
    def read(self, key: str):
        for item in reversed(self._store):
            if item["key"] == key:
                return item["value"]
        return None
    
    def all(self):
        return list(self._store)
    
    def clear(self):
        self._store = []


class BmAddAgent:
    """
    Core agent class in the BmAdd framework.
    Each agent has: identity, role, tools, memory, and an execution loop.
    """
    def __init__(self, name: str, role: str, color: str, emoji: str, tools: list = None):
        self.name = name
        self.role = role
        self.color = color
        self.emoji = emoji
        self.tools: list[BmAddTool] = tools or []
        self.memory = BmAddMemory()
        self.status = "idle"   # idle | running | done | error
        self.output = None
        self.error_msg = None
        self.start_time = None
        self.end_time = None
        self.token_count = 0
    
    def elapsed(self) -> str:
        if not self.start_time:
            return "—"
        end = self.end_time or time.time()
        return f"{end - self.start_time:.1f}s"
    
    def add_tool(self, tool: BmAddTool):
        self.tools.append(tool)
        return self
    
    def _build_system_prompt(self) -> str:
        tool_list = "\n".join([f"  - {t.name}: {t.description}" for t in self.tools])
        return f"""You are {self.name}, an expert AI agent in the BmAdd framework.
Your specialized role: {self.role}

Available tools you can conceptually use:
{tool_list}

CRITICAL RULES:
1. Always respond with VALID JSON only — no markdown, no preamble, no explanation outside JSON
2. Be precise, specific, and conversion-focused
3. Generate professional, publication-quality content
4. Never truncate arrays — complete every item fully"""
    
    async def run_async(self, task: str, api_key: str, callback=None) -> dict:
        """Async execution (called via sync wrapper in Streamlit)"""
        return self.run_sync(task, api_key, callback)
    
    def run_sync(self, task: str, api_key: str, callback=None) -> dict:
        """Synchronous execution"""
        self.status = "running"
        self.start_time = time.time()
        self.memory.write("task", task)
        
        if callback:
            callback(self.name, "running", f"Starting: {task[:60]}...")
        
        try:
            import urllib.request
            import json as _json
            
            payload = _json.dumps({
                "model": "gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": self._build_system_prompt()},
                    {"role": "user", "content": task}
                ],
                "temperature": 0.8,
                "max_tokens": 3500
            }).encode("utf-8")
            
            req = urllib.request.Request(
                "https://api.openai.com/v1/chat/completions",
                data=payload,
                headers={
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {api_key}"
                },
                method="POST"
            )
            
            with urllib.request.urlopen(req, timeout=45) as resp:
                result = _json.loads(resp.read().decode("utf-8"))
                raw = result["choices"][0]["message"]["content"]
                self.token_count = result.get("usage", {}).get("total_tokens", 0)
            
            # Parse JSON — try multiple strategies
            raw_clean = raw.strip()
            # Remove markdown fences
            raw_clean = re.sub(r'^```(?:json)?\s*', '', raw_clean, flags=re.MULTILINE)
            raw_clean = re.sub(r'\s*```\s*$', '', raw_clean, flags=re.MULTILINE)
            raw_clean = raw_clean.strip()
            
            # Extract JSON object/array
            json_match = re.search(r'\{.*\}', raw_clean, re.DOTALL)
            if json_match:
                self.output = _json.loads(json_match.group())
            else:
                self.output = _json.loads(raw_clean)
            
            self.memory.write("output", self.output)
            self.status = "done"
            self.end_time = time.time()
            
            if callback:
                callback(self.name, "done", f"Completed in {self.elapsed()} · {self.token_count} tokens")
            
            return self.output
            
        except urllib.error.HTTPError as e:
            err_body = e.read().decode("utf-8")
            self.error_msg = f"HTTP {e.code}: {err_body[:200]}"
            self.status = "error"
            self.end_time = time.time()
            if callback:
                callback(self.name, "error", f"HTTP Error {e.code}: {err_body[:100]}")
            return None
        except Exception as e:
            self.error_msg = str(e)
            self.status = "error"
            self.end_time = time.time()
            if callback:
                callback(self.name, "error", f"Error: {str(e)[:100]}")
            return None


class BmAddOrchestrator:
    """
    Orchestrates multiple agents, manages execution order,
    passes context between agents, and tracks overall pipeline state.
    """
    def __init__(self):
        self.agents: dict[str, BmAddAgent] = {}
        self.pipeline: list[str] = []
        self.context: dict = {}
        self.log: list[dict] = []
        self.total_tokens = 0
        self.start_time = None
        self.status = "idle"
    
    def register(self, agent: BmAddAgent):
        self.agents[agent.name] = agent
        return self
    
    def set_pipeline(self, order: list[str]):
        self.pipeline = order
        return self
    
    def _add_log(self, agent: str, status: str, msg: str):
        entry = {
            "agent": agent,
            "status": status,
            "msg": msg,
            "ts": datetime.now().strftime("%H:%M:%S")
        }
        self.log.append(entry)
        # Also write to session state log in real time
        if "agent_log" in st.session_state:
            st.session_state.agent_log.append(entry)
    
    def run_agent(self, name: str, task: str, api_key: str) -> dict:
        if name not in self.agents:
            return None
        agent = self.agents[name]
        # Update session state for live UI
        if "agent_statuses" in st.session_state:
            st.session_state.agent_statuses[name] = "running"
        
        result = agent.run_sync(task, api_key, callback=self._add_log)
        self.total_tokens += agent.token_count
        self.context[name] = result
        
        if "agent_statuses" in st.session_state:
            st.session_state.agent_statuses[name] = agent.status
        
        return result
    
    def get_context(self, *agent_names) -> dict:
        return {k: self.context.get(k) for k in agent_names}
    
    def reset(self):
        self.log = []
        self.context = {}
        self.total_tokens = 0
        self.start_time = None
        self.status = "idle"
        for agent in self.agents.values():
            agent.status = "idle"
            agent.output = None
            agent.error_msg = None
            agent.memory.clear()


# ─────────────────────────────────────────────────────────────────
# REGISTER ALL AGENTS
# ─────────────────────────────────────────────────────────────────

def build_orchestrator() -> BmAddOrchestrator:
    orch = BmAddOrchestrator()
    
    content_agent = BmAddAgent(
        name="ContentAgent",
        role="Elite conversion copywriter — writes high-converting landing page headlines, body copy, CTAs, social proof, and feature descriptions",
        color="var(--agent-content)",
        emoji="✍️",
        tools=[
            BmAddTool("write_headline", "Craft power-word-rich headlines"),
            BmAddTool("write_subheadline", "Write value-proposition subheadlines"),
            BmAddTool("write_body_copy", "Create emotionally resonant body text"),
            BmAddTool("write_cta", "Generate high-conversion CTA copy"),
            BmAddTool("write_features", "Develop benefit-focused feature bullets"),
            BmAddTool("write_testimonials", "Craft credible social proof testimonials"),
        ]
    )
    
    pptx_agent = BmAddAgent(
        name="PptxAgent",
        role="PowerPoint architect — designs slide decks with strategic narrative structure, compelling slides, and presenter notes",
        color="var(--agent-pptx)",
        emoji="📊",
        tools=[
            BmAddTool("plan_narrative", "Structure slide flow and story arc"),
            BmAddTool("write_title_slide", "Design cover and title slides"),
            BmAddTool("write_content_slides", "Create problem/solution/feature slides"),
            BmAddTool("write_data_slides", "Build stats and metrics slides"),
            BmAddTool("write_closing_slide", "Craft CTA and closing slides"),
            BmAddTool("add_speaker_notes", "Generate presenter speaker notes"),
        ]
    )
    
    pdf_agent = BmAddAgent(
        name="PdfAgent",
        role="Multi-document specialist — creates executive briefs, sales proposals, product overviews, and onepagers with professional structure",
        color="var(--agent-pdf)",
        emoji="📄",
        tools=[
            BmAddTool("write_executive_brief", "Create C-suite executive summary"),
            BmAddTool("write_sales_proposal", "Build persuasive sales proposals"),
            BmAddTool("write_product_overview", "Develop detailed product briefs"),
            BmAddTool("write_onepager", "Generate punchy one-page summaries"),
            BmAddTool("write_case_study", "Draft customer success case studies"),
            BmAddTool("write_datasheet", "Produce technical datasheets"),
        ]
    )
    
    seo_agent = BmAddAgent(
        name="SeoAgent",
        role="SEO strategist — generates meta tags, schema markup, keyword strategy, and OpenGraph data for maximum search visibility",
        color="var(--agent-seo)",
        emoji="🔍",
        tools=[
            BmAddTool("generate_meta_tags", "Write title and description meta tags"),
            BmAddTool("generate_og_tags", "Create OpenGraph social preview tags"),
            BmAddTool("write_schema_json_ld", "Generate JSON-LD structured data"),
            BmAddTool("keyword_strategy", "Develop primary and LSI keyword list"),
            BmAddTool("generate_sitemap_entry", "Create sitemap XML entries"),
        ]
    )
    
    design_agent = BmAddAgent(
        name="DesignAgent",
        role="Visual design strategist — recommends color palettes, typography pairings, layout systems, and brand identity for maximum impact",
        color="var(--agent-design)",
        emoji="🎨",
        tools=[
            BmAddTool("suggest_color_palette", "Generate brand-appropriate color system"),
            BmAddTool("pick_font_pairing", "Recommend distinctive font combinations"),
            BmAddTool("plan_layout_system", "Design grid and spacing system"),
            BmAddTool("design_component_library", "Define reusable UI component styles"),
        ]
    )
    
    qa_agent = BmAddAgent(
        name="QaAgent",
        role="Quality assurance agent — reviews all generated content for consistency, completeness, brand voice, and conversion optimization",
        color="var(--agent-qa)",
        emoji="✅",
        tools=[
            BmAddTool("check_consistency", "Verify brand voice across all content"),
            BmAddTool("check_completeness", "Ensure all required fields are populated"),
            BmAddTool("suggest_improvements", "Identify conversion optimization gaps"),
            BmAddTool("generate_ab_variants", "Create A/B test variants for CTAs"),
        ]
    )
    
    copy_agent = BmAddAgent(
        name="CopyAgent",
        role="Supplementary copywriter — generates email sequences, ad copy, social media posts, and ancillary marketing materials",
        color="var(--agent-copy)",
        emoji="📢",
        tools=[
            BmAddTool("write_email_sequence", "Create drip email sequences"),
            BmAddTool("write_ad_copy", "Generate PPC/social ad variations"),
            BmAddTool("write_social_posts", "Draft platform-specific social content"),
            BmAddTool("write_press_release", "Compose formal press releases"),
        ]
    )
    
    for agent in [content_agent, pptx_agent, pdf_agent, seo_agent, design_agent, qa_agent, copy_agent]:
        orch.register(agent)
    
    return orch


# ─────────────────────────────────────────────────────────────────
# PPTX GENERATOR (python-pptx)
# ─────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_color: str) -> tuple:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def generate_pptx(slides_data: list, brand_hex: str = "#4d6fff", accent_hex: str = "#ff3355", 
                  company_name: str = "Brand", font_header: str = "Calibri") -> bytes:
    """Generate a full branded .pptx file using python-pptx"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    
    brand_rgb = RGBColor(*hex_to_rgb(brand_hex))
    accent_rgb = RGBColor(*hex_to_rgb(accent_hex))
    dark_rgb = RGBColor(8, 10, 22)
    white_rgb = RGBColor(228, 231, 240)
    muted_rgb = RGBColor(82, 88, 112)
    light_bg_rgb = RGBColor(248, 249, 255)
    
    prs = Presentation()
    prs.core_properties.author = "PageForge PRO"
    prs.core_properties.title = company_name
    
    # Set slide size (16:9 widescreen)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    W = prs.slide_width
    H = prs.slide_height
    
    def add_bg_rect(slide, x, y, w, h, color_rgb, transparency=0):
        from pptx.util import Emu
        from pptx.oxml.ns import qn
        shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
        shape.line.fill.background()
        shape.line.width = 0
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb
        if transparency > 0:
            from lxml import etree
            solidFill = fill._xPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    alpha_val = int((1 - transparency/100) * 100000)
                    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                    alpha.set('val', str(alpha_val))
        return shape
    
    def add_text_box(slide, text, x, y, w, h, 
                     font_size=18, bold=False, color=white_rgb, 
                     align=PP_ALIGN.LEFT, font_face="Calibri", italic=False,
                     wrap=True):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        txBox.text_frame.word_wrap = wrap
        p = txBox.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_face
        run.font.color.rgb = color
        return txBox
    
    def add_accent_bar(slide, color_rgb, alpha_y=H - Inches(0.18), height=Inches(0.18)):
        shape = add_bg_rect(slide, 0, alpha_y, W, height, color_rgb)
        return shape
    
    for idx, slide_data in enumerate(slides_data):
        is_cover = idx == 0
        is_closing = idx == len(slides_data) - 1
        
        if is_cover:
            # ── COVER SLIDE: dark bg + brand gradient block ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # Full dark background
            add_bg_rect(slide, 0, 0, W, H, dark_rgb)
            # Left brand bar
            add_bg_rect(slide, 0, 0, Inches(0.35), H, brand_rgb)
            # Title
            title_txt = slide_data.get("title", "Presentation Title")
            add_text_box(slide, title_txt,
                         Inches(0.7), Inches(1.8), Inches(11.5), Inches(2.5),
                         font_size=44, bold=True, color=white_rgb, 
                         align=PP_ALIGN.LEFT, font_face=font_header)
            # Subtitle
            subtitle = slide_data.get("subtitle", "") or slide_data.get("body", "")
            if subtitle:
                add_text_box(slide, subtitle,
                             Inches(0.7), Inches(4.4), Inches(10), Inches(1.2),
                             font_size=16, bold=False, color=muted_rgb,
                             align=PP_ALIGN.LEFT)
            # Company name
            add_text_box(slide, company_name,
                         Inches(0.7), Inches(6.5), Inches(5), Inches(0.6),
                         font_size=13, bold=True, color=brand_rgb,
                         align=PP_ALIGN.LEFT, font_face=font_header)
            # Date
            date_str = datetime.now().strftime("%B %Y")
            add_text_box(slide, date_str,
                         Inches(10), Inches(6.5), Inches(3), Inches(0.5),
                         font_size=11, color=muted_rgb, align=PP_ALIGN.RIGHT)
            # Bottom accent
            add_accent_bar(slide, accent_rgb)
            # Slide number (skip on cover)
            
        elif is_closing:
            # ── CLOSING SLIDE: brand bg, white text ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg_rect(slide, 0, 0, W, H, brand_rgb)
            add_bg_rect(slide, 0, 0, Inches(0.35), H, dark_rgb)
            
            title_txt = slide_data.get("title", "Thank You")
            add_text_box(slide, title_txt,
                         Inches(0.7), Inches(2.2), Inches(11), Inches(2.0),
                         font_size=48, bold=True, color=white_rgb,
                         align=PP_ALIGN.LEFT, font_face=font_header)
            
            cta = slide_data.get("cta") or slide_data.get("body", "")
            if cta:
                add_text_box(slide, cta,
                             Inches(0.7), Inches(4.3), Inches(10), Inches(1.5),
                             font_size=18, color=RGBColor(200, 215, 255),
                             align=PP_ALIGN.LEFT)
            add_text_box(slide, company_name,
                         Inches(0.7), Inches(6.5), Inches(8), Inches(0.6),
                         font_size=13, bold=True, color=white_rgb,
                         align=PP_ALIGN.LEFT)
            add_accent_bar(slide, dark_rgb)
            
        elif slide_data.get("layout") == "stats":
            # ── STATS SLIDE ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg_rect(slide, 0, 0, W, H, light_bg_rgb)
            add_bg_rect(slide, 0, 0, Inches(0.35), H, brand_rgb)
            
            title_txt = slide_data.get("title", "Key Metrics")
            add_text_box(slide, title_txt,
                         Inches(0.7), Inches(0.35), Inches(11), Inches(0.85),
                         font_size=28, bold=True, color=dark_rgb,
                         align=PP_ALIGN.LEFT, font_face=font_header)
            
            stats = slide_data.get("stats", [])
            if stats:
                stat_w = Inches(11.2 / max(len(stats), 1))
                for i, stat in enumerate(stats[:4]):
                    sx = Inches(0.7) + i * stat_w
                    # Stat box bg
                    add_bg_rect(slide, sx, Inches(1.8), stat_w - Inches(0.15), Inches(3.2),
                                RGBColor(240, 242, 255))
                    # Number
                    add_text_box(slide, stat.get("num", "—"),
                                 sx + Inches(0.1), Inches(2.0), stat_w - Inches(0.3), Inches(1.5),
                                 font_size=40, bold=True, color=brand_rgb,
                                 align=PP_ALIGN.CENTER, font_face=font_header)
                    # Label
                    add_text_box(slide, stat.get("label", ""),
                                 sx + Inches(0.1), Inches(3.6), stat_w - Inches(0.3), Inches(0.8),
                                 font_size=12, color=muted_rgb, align=PP_ALIGN.CENTER)
            
            body = slide_data.get("body", "")
            if body:
                add_text_box(slide, body,
                             Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.2),
                             font_size=13, color=muted_rgb, align=PP_ALIGN.LEFT)
            add_accent_bar(slide, brand_rgb)
            add_text_box(slide, str(idx + 1), Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.4),
                         font_size=9, color=muted_rgb, align=PP_ALIGN.RIGHT)
            
        else:
            # ── STANDARD CONTENT SLIDE ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg_rect(slide, 0, 0, W, H, light_bg_rgb)
            add_bg_rect(slide, 0, 0, Inches(0.35), H, brand_rgb)
            
            title_txt = slide_data.get("title", "")
            add_text_box(slide, title_txt,
                         Inches(0.7), Inches(0.35), Inches(11), Inches(0.85),
                         font_size=28, bold=True, color=dark_rgb,
                         align=PP_ALIGN.LEFT, font_face=font_header)
            
            # Separator line (thin rect, not accent)
            add_bg_rect(slide, Inches(0.7), Inches(1.25), Inches(11.3), Pt(1.5), muted_rgb)
            
            bullets = slide_data.get("bullets", [])
            body = slide_data.get("body", "")
            notes = slide_data.get("speaker_notes", "")
            icon = slide_data.get("icon", "")
            
            if icon:
                add_text_box(slide, icon,
                             Inches(0.7), Inches(1.5), Inches(0.8), Inches(0.8),
                             font_size=28, color=brand_rgb, align=PP_ALIGN.LEFT)
            
            content_x = Inches(0.7)
            content_y = Inches(1.55)
            
            if bullets:
                for bi, bullet in enumerate(bullets[:6]):
                    by = content_y + bi * Inches(0.82)
                    if by + Inches(0.7) > H - Inches(0.5):
                        break
                    # Bullet dot
                    dot = slide.shapes.add_shape(9, content_x, by + Inches(0.18), 
                                                  Inches(0.1), Inches(0.1))  # oval
                    dot.fill.solid()
                    dot.fill.fore_color.rgb = brand_rgb
                    dot.line.fill.background()
                    # Bullet text
                    add_text_box(slide, bullet,
                                 content_x + Inches(0.2), by, Inches(11.2), Inches(0.75),
                                 font_size=15, color=dark_rgb, align=PP_ALIGN.LEFT)
            elif body:
                # Multi-paragraph body
                add_text_box(slide, body,
                             content_x, content_y, Inches(11.3), Inches(5.0),
                             font_size=15, color=dark_rgb, align=PP_ALIGN.LEFT)
            
            # Speaker notes
            if notes:
                notes_slide = slide.notes_slide
                tf = notes_slide.notes_text_frame
                tf.text = notes
            
            # Bottom bar
            add_accent_bar(slide, brand_rgb)
            # Company name footer
            add_text_box(slide, company_name,
                         Inches(0.7), Inches(6.9), Inches(5), Inches(0.45),
                         font_size=8, color=muted_rgb, align=PP_ALIGN.LEFT)
            # Slide number
            add_text_box(slide, str(idx + 1), Inches(12.5), Inches(6.9), 
                         Inches(0.5), Inches(0.45),
                         font_size=9, color=muted_rgb, align=PP_ALIGN.RIGHT)
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────
# PDF GENERATORS (reportlab)
# ─────────────────────────────────────────────────────────────────

def generate_pdf_document(doc_data: dict, brand_hex: str = "#4d6fff", company_name: str = "Brand") -> bytes:
    """Generate a single polished PDF document via reportlab"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm, cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table, 
                                     TableStyle, HRFlowable, PageBreak, KeepTogether)
    from reportlab.platypus import BaseDocTemplate, Frame, PageTemplate
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    
    br, bg, bb = hex_to_rgb(brand_hex)
    brand_color = colors.Color(br/255, bg/255, bb/255)
    brand_light = colors.Color(br/255, bg/255, bb/255, alpha=0.08)
    dark_color = colors.Color(0.06, 0.07, 0.12)
    muted_color = colors.Color(0.42, 0.44, 0.52)
    border_color = colors.Color(0.88, 0.89, 0.93)
    
    buf = io.BytesIO()
    
    page_w, page_h = A4
    margin = 2.2 * cm
    
    styles = getSampleStyleSheet()
    
    # Custom styles
    doc_type_style = ParagraphStyle("DocType",
        fontName="Helvetica-Bold", fontSize=8, textColor=brand_color,
        spaceAfter=4, letterSpacing=2.5, leading=12)
    
    title_style = ParagraphStyle("Title",
        fontName="Helvetica-Bold", fontSize=28, textColor=dark_color,
        spaceAfter=8, leading=34)
    
    subtitle_style = ParagraphStyle("Subtitle",
        fontName="Helvetica", fontSize=13, textColor=muted_color,
        spaceAfter=16, leading=20)
    
    h2_style = ParagraphStyle("H2",
        fontName="Helvetica-Bold", fontSize=16, textColor=dark_color,
        spaceBefore=20, spaceAfter=8, leading=22)
    
    h3_style = ParagraphStyle("H3",
        fontName="Helvetica-Bold", fontSize=12, textColor=dark_color,
        spaceBefore=12, spaceAfter=6, leading=18)
    
    body_style = ParagraphStyle("Body",
        fontName="Helvetica", fontSize=10.5, textColor=dark_color,
        spaceAfter=10, leading=17)
    
    highlight_style = ParagraphStyle("Highlight",
        fontName="Helvetica-Oblique", fontSize=11, textColor=dark_color,
        spaceAfter=10, leading=18,
        leftIndent=14, borderColor=brand_color, borderWidth=3,
        borderPadding=(8, 8, 8, 12))
    
    bullet_style = ParagraphStyle("Bullet",
        fontName="Helvetica", fontSize=10.5, textColor=dark_color,
        spaceAfter=6, leading=16, leftIndent=16, bulletIndent=0)
    
    caption_style = ParagraphStyle("Caption",
        fontName="Helvetica", fontSize=8.5, textColor=muted_color,
        spaceAfter=4, leading=13, alignment=TA_RIGHT)
    
    story = []
    
    # ── HEADER ──
    # Brand bar (using table to simulate colored header)
    header_data = [[Paragraph(f'<font color="white"><b>{company_name.upper()}</b></font>', 
                              ParagraphStyle("BrandHeader", fontName="Helvetica-Bold",
                                             fontSize=10, textColor=colors.white, leading=14))]]
    header_table = Table(header_data, colWidths=[page_w - 2*margin])
    header_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), brand_color),
        ("TOPPADDING", (0,0), (-1,-1), 10),
        ("BOTTOMPADDING", (0,0), (-1,-1), 10),
        ("LEFTPADDING", (0,0), (-1,-1), 14),
        ("RIGHTPADDING", (0,0), (-1,-1), 14),
    ]))
    story.append(header_table)
    story.append(Spacer(1, 16))
    
    # ── DOCUMENT META ──
    doc_type = doc_data.get("type", "Business Document")
    story.append(Paragraph(doc_type.upper(), doc_type_style))
    story.append(Paragraph(doc_data.get("title", "Untitled Document"), title_style))
    
    if doc_data.get("subtitle"):
        story.append(Paragraph(doc_data["subtitle"], subtitle_style))
    
    # Meta info table
    meta_items = [
        ["Prepared by:", doc_data.get("author", company_name)],
        ["Date:", datetime.now().strftime("%B %d, %Y")],
        ["Classification:", doc_data.get("classification", "Confidential")],
    ]
    meta_table = Table(meta_items, colWidths=[3*cm, 12*cm])
    meta_table.setStyle(TableStyle([
        ("FONTNAME", (0,0), (0,-1), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("TEXTCOLOR", (0,0), (0,-1), muted_color),
        ("TEXTCOLOR", (1,0), (1,-1), dark_color),
        ("TOPPADDING", (0,0), (-1,-1), 3),
        ("BOTTOMPADDING", (0,0), (-1,-1), 3),
        ("LEFTPADDING", (0,0), (-1,-1), 0),
    ]))
    story.append(meta_table)
    story.append(Spacer(1, 6))
    story.append(HRFlowable(width="100%", thickness=2, color=brand_color, spaceAfter=18))
    
    # ── SECTIONS ──
    for section in doc_data.get("sections", []):
        sec_title = section.get("title", "")
        if sec_title:
            story.append(Paragraph(sec_title, h2_style))
            story.append(HRFlowable(width="100%", thickness=0.5, color=border_color, spaceAfter=10))
        
        if section.get("intro"):
            story.append(Paragraph(section["intro"], body_style))
        
        if section.get("highlight"):
            # Highlighted quote block
            highlight_data = [[Paragraph(section["highlight"], 
                                         ParagraphStyle("HLInner", fontName="Helvetica-Oblique",
                                                         fontSize=11, textColor=dark_color, leading=18))]]
            hl_table = Table(highlight_data, colWidths=[page_w - 2*margin])
            hl_table.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,-1), brand_light),
                ("LEFTPADDING", (0,0), (-1,-1), 16),
                ("RIGHTPADDING", (0,0), (-1,-1), 14),
                ("TOPPADDING", (0,0), (-1,-1), 10),
                ("BOTTOMPADDING", (0,0), (-1,-1), 10),
                ("LINEAFTER", (0,0), (-1,-1), 0),
            ]))
            story.append(hl_table)
            story.append(Spacer(1, 10))
        
        if section.get("body"):
            story.append(Paragraph(section["body"], body_style))
        
        if section.get("bullets"):
            for bullet in section["bullets"]:
                story.append(Paragraph(f"&#8226;&nbsp;&nbsp;{bullet}", bullet_style))
            story.append(Spacer(1, 4))
        
        if section.get("stats"):
            stats = section["stats"]
            n = min(len(stats), 4)
            stat_data = [[
                Paragraph(f'<b>{s["num"]}</b>', ParagraphStyle("StatNum", 
                    fontName="Helvetica-Bold", fontSize=22, textColor=brand_color,
                    alignment=TA_CENTER, leading=28))
                for s in stats[:n]
            ], [
                Paragraph(s.get("label", ""), ParagraphStyle("StatLabel",
                    fontName="Helvetica", fontSize=8.5, textColor=muted_color,
                    alignment=TA_CENTER, leading=13))
                for s in stats[:n]
            ]]
            col_w = (page_w - 2*margin) / n
            stat_table = Table(stat_data, colWidths=[col_w] * n)
            stat_table.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,-1), colors.Color(0.97, 0.97, 1.0)),
                ("TOPPADDING", (0,0), (-1,-1), 14),
                ("BOTTOMPADDING", (0,0), (-1,-1), 14),
                ("LINEBELOW", (0,0), (-1,0), 0.5, border_color),
                ("LINEBEFORE", (1,0), (-1,-1), 0.5, border_color),
                ("INNERGRID", (0,0), (-1,-1), 0.5, border_color),
                ("BOX", (0,0), (-1,-1), 0.5, border_color),
            ]))
            story.append(stat_table)
            story.append(Spacer(1, 12))
        
        if section.get("table"):
            tbl_data = section["table"]
            if tbl_data:
                col_w = (page_w - 2*margin) / max(len(tbl_data[0]), 1)
                tbl = Table(tbl_data, colWidths=[col_w] * len(tbl_data[0]))
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (0,0), (-1,0), brand_color),
                    ("TEXTCOLOR", (0,0), (-1,0), colors.white),
                    ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
                    ("FONTSIZE", (0,0), (-1,-1), 9.5),
                    ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.Color(0.97,0.97,1)]),
                    ("INNERGRID", (0,0), (-1,-1), 0.5, border_color),
                    ("BOX", (0,0), (-1,-1), 0.5, border_color),
                    ("TOPPADDING", (0,0), (-1,-1), 6),
                    ("BOTTOMPADDING", (0,0), (-1,-1), 6),
                    ("LEFTPADDING", (0,0), (-1,-1), 8),
                ]))
                story.append(tbl)
                story.append(Spacer(1, 12))
        
        story.append(Spacer(1, 6))
    
    # ── FOOTER INFO ──
    story.append(Spacer(1, 20))
    story.append(HRFlowable(width="100%", thickness=0.5, color=border_color, spaceAfter=8))
    story.append(Paragraph(
        f"{doc_data.get('title','Document')} &nbsp;·&nbsp; {company_name} &nbsp;·&nbsp; Confidential &nbsp;·&nbsp; {datetime.now().strftime('%Y')}",
        caption_style
    ))
    
    # Build PDF
    def on_first_page(canvas, doc):
        pass
    
    def on_later_pages(canvas, doc):
        canvas.saveState()
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(muted_color)
        canvas.drawRightString(page_w - margin, 1.2*cm, f"Page {doc.page}")
        canvas.drawString(margin, 1.2*cm, company_name)
        canvas.restoreState()
    
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=margin, rightMargin=margin,
        topMargin=margin, bottomMargin=1.8*cm,
        title=doc_data.get("title", "Document"),
        author=company_name
    )
    doc.build(story, onFirstPage=on_first_page, onLaterPages=on_later_pages)
    buf.seek(0)
    return buf.getvalue()


def generate_all_pdfs(pdf_docs: list, brand_hex: str, company_name: str) -> bytes:
    """Generate all PDFs and return as a ZIP archive"""
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, doc in enumerate(pdf_docs):
            pdf_bytes = generate_pdf_document(doc, brand_hex, company_name)
            safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', doc.get("type", f"doc_{i}"))
            zf.writestr(f"{safe_name}.pdf", pdf_bytes)
    zip_buf.seek(0)
    return zip_buf.getvalue()


# ─────────────────────────────────────────────────────────────────
# LANDING PAGE HTML GENERATOR
# ─────────────────────────────────────────────────────────────────

def hex_to_rgb_str(hex_color: str) -> str:
    r, g, b = hex_to_rgb(hex_color)
    return f"{r}, {g}, {b}"

def generate_html_page(data: dict, settings: dict) -> str:
    brand = settings.get("brand_color", "#4d6fff")
    accent = settings.get("accent_color", "#ff3355")
    bg = settings.get("bg_color", "#060810")
    text_color = settings.get("text_color", "#e4e7f0")
    font_pair = settings.get("font_pair", "Modern")
    hero_layout = settings.get("hero_layout", "Centered")
    animation = settings.get("animation", "Smooth")
    bg_effect = settings.get("bg_effect", "Gradient Mesh")
    
    brand_rgb = hex_to_rgb_str(brand)
    accent_rgb = hex_to_rgb_str(accent)
    bg_rgb = hex_to_rgb(bg)
    is_dark = sum(bg_rgb) / 3 < 128
    
    fonts = {
        "Modern": ("https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=DM+Sans:wght@300;400;500&display=swap", "'Syne', sans-serif", "'DM Sans', sans-serif"),
        "Editorial": ("https://fonts.googleapis.com/css2?family=Playfair+Display:ital,wght@0,700;1,400&family=Lato:wght@300;400&display=swap", "'Playfair Display', serif", "'Lato', sans-serif"),
        "Technical": ("https://fonts.googleapis.com/css2?family=IBM+Plex+Mono:wght@400;600&family=IBM+Plex+Sans:wght@300;400;600&display=swap", "'IBM Plex Mono', monospace", "'IBM Plex Sans', sans-serif"),
        "Bold": ("https://fonts.googleapis.com/css2?family=Bebas+Neue&family=Inter:wght@300;400;500&display=swap", "'Bebas Neue', display", "'Inter', sans-serif"),
        "Elegant": ("https://fonts.googleapis.com/css2?family=Cormorant+Garamond:ital,wght@0,600;1,400&family=Josefin+Sans:wght@300;400;600&display=swap", "'Cormorant Garamond', serif", "'Josefin Sans', sans-serif"),
        "Futuristic": ("https://fonts.googleapis.com/css2?family=Orbitron:wght@400;700;900&family=Rajdhani:wght@300;400;600&display=swap", "'Orbitron', sans-serif", "'Rajdhani', sans-serif"),
    }
    font_import, font_heading, font_body = fonts.get(font_pair, fonts["Modern"])
    
    bg_css = {
        "Gradient Mesh": f"background: radial-gradient(ellipse at 0% 0%, rgba({brand_rgb}, 0.14) 0%, transparent 55%), radial-gradient(ellipse at 100% 100%, rgba({accent_rgb}, 0.10) 0%, transparent 55%), {bg};",
        "Dots": f"background-color: {bg}; background-image: radial-gradient(rgba({brand_rgb}, 0.12) 1px, transparent 1px); background-size: 26px 26px;",
        "Lines": f"background-color: {bg}; background-image: repeating-linear-gradient(0deg, rgba({brand_rgb}, 0.05) 0px, rgba({brand_rgb}, 0.05) 1px, transparent 1px, transparent 38px);",
        "Solid": f"background-color: {bg};",
    }.get(bg_effect, f"background-color: {bg};")
    
    anim_css = {
        "Smooth": "@keyframes fadeUp { from { opacity:0; transform:translateY(28px); } to { opacity:1; transform:translateY(0); } } @keyframes fadeIn { from { opacity:0; } to { opacity:1; } } .au { animation: fadeUp 0.65s ease forwards; } .ai { animation: fadeIn 0.55s ease forwards; } .d1{animation-delay:.1s;opacity:0} .d2{animation-delay:.2s;opacity:0} .d3{animation-delay:.3s;opacity:0} .d4{animation-delay:.4s;opacity:0} .d5{animation-delay:.5s;opacity:0} .d6{animation-delay:.6s;opacity:0}",
        "Dramatic": "@keyframes zoomIn { from { opacity:0; transform:scale(0.88); } to { opacity:1; transform:scale(1); } } .au,.ai { animation: zoomIn 0.85s cubic-bezier(0.16,1,0.3,1) forwards; } .d1{animation-delay:.15s;opacity:0} .d2{animation-delay:.3s;opacity:0} .d3{animation-delay:.45s;opacity:0} .d4{animation-delay:.6s;opacity:0} .d5{animation-delay:.75s;opacity:0} .d6{animation-delay:.9s;opacity:0}",
        "None": ".au,.ai{} .d1,.d2,.d3,.d4,.d5,.d6{}"
    }.get(animation, "")
    
    surface = "rgba(255,255,255,0.05)" if is_dark else "rgba(0,0,0,0.04)"
    border_c = "rgba(255,255,255,0.09)" if is_dark else "rgba(0,0,0,0.09)"
    muted_c = "rgba(255,255,255,0.45)" if is_dark else "rgba(0,0,0,0.42)"
    
    nav_links_html = "".join([f'<a href="#{l.lower()}">{l}</a>' for l in data.get("nav_links", ["Features","Testimonials","FAQ"])])
    
    feat_html = ""
    for i, f in enumerate(data.get("features", [])):
        feat_html += f'<div class="fc au d{min(i+1,6)}" id="f{i}"><span class="fi">{f.get("icon","⚡")}</span><h3>{f.get("title","")}</h3><p>{f.get("desc","")}</p></div>'
    
    test_html = ""
    for t in data.get("testimonials", []):
        stars = "★" * t.get("rating", 5)
        test_html += f'<div class="tc"><div class="ts">{stars}</div><p>"{t.get("text","")}"</p><div class="ta"><div class="tav">{(t.get("name","A"))[0]}</div><div><b>{t.get("name","")}</b><small>{t.get("role","")}</small></div></div></div>'
    
    faq_html = ""
    for i, item in enumerate(data.get("faq", [])):
        faq_html += f'<div class="fqi"><button class="fqb" onclick="tFAQ(this)"><span>{item.get("q","")}</span><span class="fqi-icon">+</span></button><div class="fqa"><p>{item.get("a","")}</p></div></div>'
    
    stats_html = ""
    for s in data.get("stats", []):
        stats_html += f'<div class="si"><div class="sn">{s.get("num","")}</div><div class="sl">{s.get("label","")}</div></div>'
    
    seo = data.get("seo_meta", {})
    meta_desc = seo.get("meta_description", data.get("subheadline", ""))
    og_title = seo.get("og_title", data.get("headline", ""))
    schema = seo.get("schema_json_ld", "")
    
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width,initial-scale=1">
    <meta name="description" content="{meta_desc}">
    <meta property="og:title" content="{og_title}">
    <meta property="og:description" content="{meta_desc}">
    <meta name="theme-color" content="{brand}">
    <title>{data.get('headline','')} — {data.get('brand_name','Brand')}</title>
    {f'<script type="application/ld+json">{schema}</script>' if schema else ''}
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
    <link href="{font_import}" rel="stylesheet">
    <style>
        *,*::before,*::after{{box-sizing:border-box;margin:0;padding:0}}
        :root{{
            --brand:{brand}; --accent:{accent}; --bg:{bg}; --text:{text_color};
            --br:{brand_rgb}; --ar:{accent_rgb}; --surface:{surface};
            --border:{border_c}; --muted:{muted_c};
            --hf:{font_heading}; --bf:{font_body};
        }}
        html{{scroll-behavior:smooth}}
        body{{{bg_css} color:var(--text); font-family:var(--bf); line-height:1.65; -webkit-font-smoothing:antialiased}}
        {anim_css}
        .wrap{{max-width:1160px;margin:0 auto;padding:0 24px}}
        /* NAV */
        nav{{background:{'rgba(6,8,16,.9)' if is_dark else 'rgba(255,255,255,.9)'};backdrop-filter:blur(18px);-webkit-backdrop-filter:blur(18px);border-bottom:1px solid var(--border);{'position:sticky;top:0;z-index:100;' if settings.get('nav_sticky',True) else ''}}}
        .nav-i{{max-width:1160px;margin:0 auto;padding:0 24px;height:66px;display:flex;align-items:center;gap:2rem}}
        .nav-brand{{font-family:var(--hf);font-size:1.35rem;font-weight:800;color:var(--brand);flex-shrink:0;letter-spacing:0.03em}}
        .nav-links{{display:flex;gap:1.5rem;flex:1}} .nav-links a{{color:var(--muted);text-decoration:none;font-size:0.9rem;font-weight:500;transition:color .2s}} .nav-links a:hover{{color:var(--text)}}
        .nav-cta{{background:var(--brand);color:{'#fff' if not is_dark else '#000'};padding:8px 20px;border-radius:7px;text-decoration:none;font-weight:600;font-size:.875rem;transition:all .2s;flex-shrink:0}} .nav-cta:hover{{opacity:.88;transform:translateY(-1px)}}
        .nav-hb{{display:none;background:none;border:none;color:var(--text);font-size:1.4rem;cursor:pointer}}
        /* HERO */
        .hero{{padding:95px 0 75px;position:relative;overflow:hidden}}
        .hero-bg{{position:absolute;top:-180px;left:50%;transform:translateX(-50%);width:750px;height:750px;background:radial-gradient(circle,rgba(var(--br),.11) 0%,transparent 65%);pointer-events:none}}
        .h-badge{{display:inline-flex;background:rgba(var(--br),.12);border:1px solid rgba(var(--br),.3);color:var(--brand);padding:5px 16px;border-radius:100px;font-size:.8rem;font-weight:600;margin-bottom:1.4rem;letter-spacing:.03em}}
        .hero-c{{text-align:center;max-width:800px;margin:0 auto}}
        .hero-c h1{{font-family:var(--hf);font-size:clamp(2.5rem,6vw,4.5rem);font-weight:800;line-height:1.1;letter-spacing:-.02em;margin-bottom:1.2rem;background:linear-gradient(135deg,var(--text) 0%,var(--brand) 100%);-webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text}}
        .hero-sub{{font-size:1.15rem;color:var(--muted);margin-bottom:.9rem;line-height:1.55}}
        .hero-body{{font-size:.975rem;color:var(--muted);margin-bottom:1.75rem;max-width:580px;margin-left:auto;margin-right:auto;line-height:1.7}}
        .btns{{display:flex;gap:10px;justify-content:center;flex-wrap:wrap;margin-bottom:1.4rem}}
        .hero-ef{{display:flex;gap:8px;max-width:420px;margin:0 auto;background:var(--surface);border:1px solid var(--border);border-radius:9px;padding:5px}}
        .ef-input{{flex:1;background:none;border:none;outline:none;color:var(--text);font-family:var(--bf);font-size:.9rem;padding:5px 10px}} .ef-input::placeholder{{color:var(--muted)}}
        /* HERO SPLIT */
        .hero-s{{padding:75px 0 60px}} .hero-si{{display:grid;grid-template-columns:1fr 1fr;gap:3.5rem;align-items:center}}
        .hero-s h1{{font-family:var(--hf);font-size:clamp(2rem,5vw,3.5rem);font-weight:800;line-height:1.1;letter-spacing:-.02em;margin-bottom:.9rem;background:linear-gradient(135deg,var(--text) 0%,var(--brand) 100%);-webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text}}
        .mock{{background:var(--surface);border:1px solid var(--border);border-radius:14px;overflow:hidden;box-shadow:0 22px 55px rgba(0,0,0,.3)}}
        .mock-bar{{display:flex;gap:5px;padding:10px 14px;background:rgba(255,255,255,.04);border-bottom:1px solid var(--border)}} .mock-bar span{{width:9px;height:9px;border-radius:50%}} .mock-bar span:nth-child(1){{background:#ff5f57}} .mock-bar span:nth-child(2){{background:#febc2e}} .mock-bar span:nth-child(3){{background:#28c840}}
        .mock-c{{padding:18px;display:flex;flex-direction:column;gap:10px}}
        .ml{{height:9px;background:var(--border);border-radius:3px}} .ml.w{{width:100%}} .ml.m{{width:68%}} .ml.s{{width:44%}}
        .mc{{height:75px;background:linear-gradient(135deg,rgba(var(--br),.15),rgba(var(--ar),.1));border:1px solid rgba(var(--br),.2);border-radius:7px}} .mc.sm{{height:55px;flex:1}}
        .mr{{display:flex;gap:10px}}
        /* HERO MINIMAL */
        .hero-min{{padding:115px 0 75px}} .hero-eye{{font-size:.72rem;letter-spacing:.22em;text-transform:uppercase;color:var(--brand);font-weight:700;margin-bottom:.9rem;display:block}}
        .hero-min h1{{font-family:var(--hf);font-size:clamp(3rem,8vw,7rem);font-weight:800;line-height:1;letter-spacing:-.03em;margin-bottom:1.3rem}}
        .hero-div{{width:55px;height:3px;background:var(--brand);margin-bottom:1.3rem;border-radius:2px}}
        /* BUTTONS */
        .btn{{display:inline-flex;align-items:center;gap:.4rem;padding:.75rem 1.75rem;border-radius:8px;text-decoration:none;font-family:var(--bf);font-weight:600;font-size:.9rem;transition:all .2s;cursor:pointer;border:2px solid transparent;white-space:nowrap}}
        .bp{{background:var(--brand);color:{'#fff' if not is_dark else '#000'};box-shadow:0 4px 14px rgba(var(--br),.35)}} .bp:hover{{transform:translateY(-2px);box-shadow:0 8px 22px rgba(var(--br),.45);opacity:.92}}
        .bg{{border-color:var(--border);color:var(--text);background:var(--surface)}} .bg:hover{{border-color:var(--brand);color:var(--brand)}}
        .bl{{padding:.95rem 2.2rem;font-size:1rem}}
        /* STATS */
        .stats{{padding:55px 0;border-top:1px solid var(--border);border-bottom:1px solid var(--border)}}
        .sg{{display:grid;grid-template-columns:repeat(4,1fr);gap:1.5rem;text-align:center}}
        .sn{{font-family:var(--hf);font-size:2.7rem;font-weight:800;color:var(--brand);letter-spacing:.02em;line-height:1;margin-bottom:5px}}
        .sl{{font-size:.75rem;color:var(--muted);text-transform:uppercase;letter-spacing:.12em;font-weight:600}}
        /* SECTION HEADERS */
        .sh{{text-align:center;margin-bottom:3.5rem}} .sh-eye{{display:inline-block;font-size:.7rem;letter-spacing:.2em;text-transform:uppercase;color:var(--brand);font-weight:700;background:rgba(var(--br),.1);padding:3px 10px;border-radius:100px;margin-bottom:10px}}
        .sh h2{{font-family:var(--hf);font-size:clamp(1.8rem,3.5vw,2.8rem);font-weight:800;letter-spacing:-.02em;line-height:1.2}} .sh h2 em{{font-style:normal;color:var(--brand)}}
        /* FEATURES */
        .feat{{padding:90px 0}} .fg{{display:grid;grid-template-columns:repeat(3,1fr);gap:1.4rem}}
        .fc{{background:var(--surface);border:1px solid var(--border);border-radius:11px;padding:1.75rem;transition:all .28s;position:relative;overflow:hidden}}
        .fc::before{{content:'';position:absolute;top:0;left:0;right:0;height:2px;background:linear-gradient(90deg,var(--brand),var(--accent));opacity:0;transition:opacity .28s}}
        .fc:hover{{transform:translateY(-4px);border-color:rgba(var(--br),.35);box-shadow:0 12px 28px rgba(var(--br),.12)}} .fc:hover::before{{opacity:1}}
        .fi{{font-size:1.85rem;margin-bottom:.85rem;display:block}} .fc h3{{font-weight:700;font-size:1rem;margin-bottom:7px;color:var(--text)}} .fc p{{font-size:.875rem;color:var(--muted);line-height:1.65}}
        /* TESTIMONIALS */
        .test{{padding:90px 0;background:{'rgba(255,255,255,.02)' if is_dark else 'rgba(0,0,0,.02)'}}}
        .tg{{display:grid;grid-template-columns:repeat(3,1fr);gap:1.4rem}}
        .tc{{background:{'rgba(255,255,255,.03)' if is_dark else '#fff'};border:1px solid var(--border);border-radius:11px;padding:1.6rem;display:flex;flex-direction:column;gap:.9rem}}
        .ts{{color:#fbbf24;font-size:.95rem;letter-spacing:2px}} .tc p{{font-size:.9rem;line-height:1.7;color:var(--text);flex:1;font-style:italic}}
        .ta{{display:flex;align-items:center;gap:10px}} .tav{{width:38px;height:38px;border-radius:50%;background:linear-gradient(135deg,var(--brand),var(--accent));display:flex;align-items:center;justify-content:center;font-weight:700;font-size:.9rem;color:#fff;flex-shrink:0}}
        .ta b{{font-size:.875rem;color:var(--text);display:block}} .ta small{{font-size:.75rem;color:var(--muted)}}
        /* FAQ */
        .faq{{padding:90px 0}} .fql{{max-width:680px;margin:0 auto}} .fqi{{border-bottom:1px solid var(--border)}}
        .fqb{{width:100%;background:none;border:none;cursor:pointer;display:flex;justify-content:space-between;align-items:center;padding:1.2rem 0;font-family:var(--bf);font-weight:600;font-size:.975rem;color:var(--text);text-align:left;gap:1rem}}
        .fqi-icon{{font-size:1.2rem;color:var(--brand);transition:transform .3s;flex-shrink:0}} .fqb.open .fqi-icon{{transform:rotate(45deg)}}
        .fqa{{max-height:0;overflow:hidden;transition:max-height .38s ease}} .fqa.open{{max-height:200px}} .fqa p{{padding-bottom:1.2rem;color:var(--muted);font-size:.925rem;line-height:1.7}}
        /* CTA BAND */
        .cband{{padding:75px 0;background:linear-gradient(135deg,rgba(var(--br),.12) 0%,rgba(var(--ar),.08) 100%);border-top:1px solid rgba(var(--br),.2);border-bottom:1px solid rgba(var(--br),.2)}}
        .cband-i{{display:flex;align-items:center;justify-content:space-between;gap:2rem;flex-wrap:wrap}}
        .cband h2{{font-family:var(--hf);font-size:clamp(1.5rem,3vw,2.4rem);font-weight:800;letter-spacing:-.01em;margin-bottom:.4rem}} .cband p{{color:var(--muted);font-size:.975rem}}
        .cba{{display:flex;gap:10px;flex-shrink:0;flex-wrap:wrap}}
        /* FOOTER */
        .foot{{padding:36px 0;border-top:1px solid var(--border)}}
        .foot-i{{display:flex;align-items:center;justify-content:space-between;gap:1rem;flex-wrap:wrap}}
        .foot-brand{{font-family:var(--hf);font-size:1.3rem;font-weight:800;color:var(--brand);letter-spacing:.03em;margin-bottom:3px}} .foot-tag{{font-size:.78rem;color:var(--muted)}}
        .foot-c{{font-size:.78rem;color:var(--muted)}}
        /* RESPONSIVE */
        @media(max-width:900px){{.hero-si{{grid-template-columns:1fr}}.hv{{display:none}}.fg{{grid-template-columns:1fr 1fr}}.tg{{grid-template-columns:1fr}}.sg{{grid-template-columns:2fr 2fr}}}}
        @media(max-width:640px){{.nav-links{{display:none}}.nav-links.open{{display:flex;flex-direction:column;position:absolute;top:66px;left:0;right:0;background:var(--bg);padding:1rem;border-bottom:1px solid var(--border);z-index:200}}.nav-i{{position:relative}}.nav-hb{{display:block}}.fg{{grid-template-columns:1fr}}.sg{{grid-template-columns:1fr 1fr}}.cband-i{{flex-direction:column;text-align:center}}.hero-ef{{flex-direction:column}}.btns{{flex-direction:column;align-items:center}}}}
    </style>
</head>
<body>
<nav>
    <div class="nav-i">
        <div class="nav-brand">{data.get('brand_name','Brand')}</div>
        <div class="nav-links">{nav_links_html}</div>
        <a href="{settings.get('cta_url','#')}" class="nav-cta">{data.get('cta_primary','Get Started')}</a>
        <button class="nav-hb" onclick="document.querySelector('.nav-links').classList.toggle('open')">☰</button>
    </div>
</nav>

{'<section class="hero"><div class="hero-bg"></div><div class="wrap"><div class="hero-c"><div class="h-badge au">'+data.get("social_proof","")+'</div><h1 class="au d1">'+data.get("headline","")+'</h1><p class="hero-sub au d2">'+data.get("subheadline","")+'</p><p class="hero-body au d3">'+data.get("hero_body","")+'</p><div class="btns au d4"><a href="'+settings.get("cta_url","#")+'" class="btn bp">'+data.get("cta_primary","Get Started")+'</a><a href="#features" class="btn bg">'+data.get("cta_secondary","Learn More")+' →</a></div><div class="hero-ef au d5"><input type="email" placeholder="'+data.get("email_placeholder","your@email.com")+'" class="ef-input"><a href="'+settings.get("cta_url","#")+'" class="btn bp">'+data.get("cta_primary","Start Free")+'</a></div></div></div></section>' if hero_layout == "Centered" else ""}

{'<section class="hero-s"><div class="wrap"><div class="hero-si"><div class="hero-sc"><div class="h-badge au">'+data.get("social_proof","")+'</div><h1 class="au d1">'+data.get("headline","")+'</h1><p class="hero-sub au d2">'+data.get("subheadline","")+'</p><p class="hero-body au d3">'+data.get("hero_body","")+'</p><div class="btns au d4" style="justify-content:flex-start"><a href="'+settings.get("cta_url","#")+'" class="btn bp">'+data.get("cta_primary","Get Started")+'</a><a href="#features" class="btn bg">'+data.get("cta_secondary","Learn More")+' →</a></div></div><div class="hv"><div class="mock"><div class="mock-bar"><span></span><span></span><span></span></div><div class="mock-c"><div class="ml w"></div><div class="ml m"></div><div class="ml s"></div><div class="mc"></div><div class="mr"><div class="mc sm"></div><div class="mc sm"></div></div></div></div></div></div></div></section>' if hero_layout == "Split" else ""}

{'<section class="hero-min"><div class="wrap"><span class="hero-eye au">'+data.get("tagline","")+'</span><h1 class="au d1">'+data.get("headline","")+'</h1><div class="hero-div"></div><p class="hero-sub au d2">'+data.get("subheadline","")+'</p><div class="btns au d3" style="justify-content:flex-start"><a href="'+settings.get("cta_url","#")+'" class="btn bp">'+data.get("cta_primary","Get Started")+'</a><a href="#features" class="btn bg">'+data.get("cta_secondary","Learn More")+' →</a></div></div></section>' if hero_layout == "Minimal" else ""}

<div class="stats"><div class="wrap"><div class="sg">{stats_html}</div></div></div>

<section class="feat" id="features"><div class="wrap"><div class="sh"><span class="sh-eye">Why Choose Us</span><h2>Everything you need to <em>succeed</em></h2></div><div class="fg">{feat_html}</div></div></section>

<section class="test" id="testimonials"><div class="wrap"><div class="sh"><span class="sh-eye">Social Proof</span><h2>Loved by <em>thousands</em></h2></div><div class="tg">{test_html}</div></div></section>

<section class="faq" id="faq"><div class="wrap"><div class="sh"><span class="sh-eye">FAQ</span><h2>Got <em>questions?</em></h2></div><div class="fql">{faq_html}</div></div></section>

<section class="cband"><div class="wrap"><div class="cband-i"><div><h2>{data.get('headline','Ready to start?')}</h2><p>{data.get('footer_text','')}</p></div><div class="cba"><a href="{settings.get('cta_url','#')}" class="btn bp bl">{data.get('cta_primary','Get Started')} →</a><a href="#features" class="btn bg">{data.get('cta_secondary','Learn More')}</a></div></div></div></section>

<footer class="foot"><div class="wrap"><div class="foot-i"><div><div class="foot-brand">{data.get('brand_name','Brand')}</div><div class="foot-tag">{data.get('tagline','')}</div></div><div class="foot-c">© {datetime.now().year} {data.get('brand_name','Brand')}. All rights reserved.</div></div></div></footer>

<script>
function tFAQ(btn){{btn.classList.toggle('open');btn.nextElementSibling.classList.toggle('open')}}
if('IntersectionObserver' in window){{
    const obs=new IntersectionObserver((es)=>{{es.forEach(e=>{{if(e.isIntersecting){{e.target.style.animationPlayState='running';obs.unobserve(e.target)}}}});}},{{threshold:.08}});
    document.querySelectorAll('.d1,.d2,.d3,.d4,.d5,.d6').forEach(el=>{{el.style.animationPlayState='paused';obs.observe(el);}});
}}
</script>
</body>
</html>"""


# ─────────────────────────────────────────────────────────────────
# AI GENERATION TASKS
# ─────────────────────────────────────────────────────────────────

def task_content(description, tone, industry, audience) -> str:
    return f"""Create complete landing page content for this business:
DESCRIPTION: {description}
TONE: {tone}
INDUSTRY: {industry}
TARGET AUDIENCE: {audience}

Return ONLY this exact JSON structure (no truncation, complete every field):
{{
  "headline": "Short punchy headline, max 8 words, power words",
  "subheadline": "One sentence value prop, 12-20 words, specific",
  "hero_body": "2-3 sentences emotional hook, pain points & aspirations, 55-80 words",
  "cta_primary": "2-4 word action CTA",
  "cta_secondary": "2-4 word secondary action",
  "email_placeholder": "example@domain.com",
  "social_proof": "Short credibility line, e.g. '10k+ users · 4.9★ · Used in 40 countries'",
  "brand_name": "Short brand name 1-3 words",
  "tagline": "3-6 word brand tagline",
  "footer_text": "One sentence inspirational footer",
  "nav_links": ["Features","Testimonials","Pricing","FAQ"],
  "features": [
    {{"icon":"⚡","title":"Feature Title","desc":"25-35 word compelling benefit description with specificity"}},
    {{"icon":"🎯","title":"Feature Title","desc":"25-35 word compelling benefit description with specificity"}},
    {{"icon":"🔥","title":"Feature Title","desc":"25-35 word compelling benefit description with specificity"}},
    {{"icon":"💎","title":"Feature Title","desc":"25-35 word compelling benefit description with specificity"}},
    {{"icon":"🚀","title":"Feature Title","desc":"25-35 word compelling benefit description with specificity"}},
    {{"icon":"✨","title":"Feature Title","desc":"25-35 word compelling benefit description with specificity"}}
  ],
  "testimonials": [
    {{"name":"First Last","role":"Job Title, Company","text":"35-45 word authentic testimonial with specific result or metric","rating":5}},
    {{"name":"First Last","role":"Job Title, Company","text":"35-45 word authentic testimonial with specific result or metric","rating":5}},
    {{"name":"First Last","role":"Job Title, Company","text":"35-45 word authentic testimonial with specific result or metric","rating":5}}
  ],
  "faq": [
    {{"q":"Most common objection as question?","a":"38-48 word reassuring, benefit-focused answer"}},
    {{"q":"Pricing or value concern?","a":"38-48 word reassuring, benefit-focused answer"}},
    {{"q":"How does it work or onboarding?","a":"38-48 word reassuring, benefit-focused answer"}},
    {{"q":"Results or timeline question?","a":"38-48 word reassuring, benefit-focused answer"}}
  ],
  "stats": [
    {{"num":"10k+","label":"Active Users"}},
    {{"num":"98%","label":"Satisfaction Rate"}},
    {{"num":"3.2x","label":"Average ROI"}},
    {{"num":"24/7","label":"Expert Support"}}
  ]
}}"""


def task_pptx(description, brand_name, headline, features) -> str:
    feats_str = "; ".join([f.get("title","") for f in features[:4]])
    return f"""Create a 10-slide PowerPoint presentation for: "{description}"
Brand: {brand_name}
Headline: {headline}
Key features: {feats_str}

Return ONLY this JSON (complete all slides, no truncation):
{{
  "slides": [
    {{"title":"{headline}","subtitle":"{brand_name} — Transforming the way you work","layout":"cover","speaker_notes":"Welcome and thank the audience. Set context."}},
    {{"title":"The Problem","body":"2-3 sentence description of the key pain point this solves. Be specific about the cost of the status quo.","icon":"😤","speaker_notes":"Pause here for audience reflection. Ask if they've experienced this.","layout":"content"}},
    {{"title":"Introducing {brand_name}","body":"2-3 sentence mission statement and solution overview. What makes this categorically different.","icon":"🚀","speaker_notes":"Energy up. This is the moment of relief.","layout":"content"}},
    {{"title":"How It Works","bullets":["Step 1: [specific action] in under [time]","Step 2: [specific action] that delivers [result]","Step 3: [specific outcome] you can measure","Result: [concrete transformation]"],"speaker_notes":"Walk through each step slowly. Ask if there are questions.","layout":"content"}},
    {{"title":"Key Features","bullets":["[Feature 1] — [specific benefit]","[Feature 2] — [specific benefit]","[Feature 3] — [specific benefit]","[Feature 4] — [specific benefit]"],"speaker_notes":"Tailor this to the audience's specific pain points.","layout":"content"}},
    {{"title":"Results That Speak","layout":"stats","stats":[{{"num":"10x","label":"Faster Results"}},{{"num":"68%","label":"Cost Reduction"}},{{"num":"4.9★","label":"Customer Rating"}},{{"num":"30d","label":"Time to ROI"}}],"body":"Real outcomes from real customers. Ask me for references.","speaker_notes":"Let the numbers breathe. Don't rush past this slide."}},
    {{"title":"What Our Customers Say","body":"[Customer name], [Title] at [Company]: 'Specific impactful quote about transformation and results achieved. Include a metric.' — Another quote from a second customer with a compelling result.","icon":"💬","speaker_notes":"If possible, share a live case study from a known brand.","layout":"content"}},
    {{"title":"Pricing & Plans","bullets":["Starter: [price] — [what's included]","Professional: [price] — [what's included + key upgrade]","Enterprise: Custom — [dedicated support + features]","All plans: 30-day money-back guarantee"],"icon":"💰","speaker_notes":"Address ROI payback period. Compare to status quo cost.","layout":"content"}},
    {{"title":"Why Choose {brand_name}","bullets":["[Differentiator 1] — What no competitor does","[Differentiator 2] — Unique technology or approach","[Differentiator 3] — Proven track record with evidence","[Differentiator 4] — Support and partnership model"],"speaker_notes":"Anticipate objections about alternatives. Be direct.","layout":"content"}},
    {{"title":"Start Today","cta":"Ready to transform your [key outcome]? Get started free in under 5 minutes. No credit card required. Our team is ready to help you succeed.","layout":"closing","speaker_notes":"Ask for the next step. Offer to schedule a follow-up."}}
  ]
}}"""


def task_pdfs(description, brand_name, headline, features, testimonials, stats) -> str:
    feat_list = "; ".join([f.get("title","") for f in features[:4]])
    stat_list = "; ".join([f'{s.get("num","")} {s.get("label","")}' for s in stats[:3]])
    return f"""Create 4 professional business documents for: "{description}"
Brand: {brand_name}. Headline: {headline}. 
Features: {feat_list}. Stats: {stat_list}

Return ONLY this JSON (all sections complete, no truncation):
{{
  "documents": [
    {{
      "type": "Executive Brief",
      "title": "[Brand] Executive Summary",
      "subtitle": "Strategic overview for decision-makers",
      "classification": "Confidential",
      "sections": [
        {{"title":"Executive Overview","intro":"[2-3 sentence high-level summary of the business problem, solution, and opportunity. Be specific about market size and urgency.]","highlight":"[Single most compelling insight or stat that captures the opportunity — 25-35 words]"}},
        {{"title":"The Opportunity","body":"[2-3 sentences about market size, timing, and competitive advantage. Include specific numbers.]","stats":[{{"num":"$4.2B","label":"Market Size"}},{{"num":"38%","label":"YoY Growth"}},{{"num":"2.3x","label":"ROI Multiple"}}]}},
        {{"title":"Solution Architecture","bullets":["[Key capability 1 with specific technical/functional detail]","[Key capability 2 with specific technical/functional detail]","[Key capability 3 with specific technical/functional detail]","[Key capability 4 with specific technical/functional detail]"]}},
        {{"title":"Strategic Recommendation","body":"[2-3 sentences on recommended next steps, timeline, and expected outcomes. Be action-oriented.]"}}
      ]
    }},
    {{
      "type": "Sales Proposal",
      "title": "Proposal: [Specific transformation promise]",
      "subtitle": "A partnership to achieve [key outcome]",
      "classification": "Confidential",
      "sections": [
        {{"title":"Understanding Your Challenge","body":"[2-3 sentences articulating the prospect's pain point with empathy and specificity. Mirror their language.]"}},
        {{"title":"Our Proposed Solution","intro":"[Introduction to tailored solution approach — 2 sentences]","bullets":["[Solution component 1 mapped to their need]","[Solution component 2 mapped to their need]","[Solution component 3 mapped to their need]","[Implementation timeline and milestones]"]}},
        {{"title":"Investment & ROI","highlight":"[Compelling ROI statement — how quickly they'll recoup investment and what the ongoing return looks like — 30-40 words]","stats":[{{"num":"90d","label":"Payback Period"}},{{"num":"3.2x","label":"Projected ROI"}},{{"num":"$280K","label":"Annual Savings"}}]}},
        {{"title":"Why Partner With Us","bullets":["[Proof point 1: specific customer result]","[Proof point 2: specific customer result]","[Differentiator: what makes you unique]","[Risk reduction: guarantee or support model]"]}},
        {{"title":"Proposed Next Steps","body":"[Clear action plan: meeting, pilot, timeline. Create urgency without pressure.]"}}
      ]
    }},
    {{
      "type": "Product Overview",
      "title": "[Brand] Product Deep-Dive",
      "subtitle": "Technical and functional capabilities guide",
      "classification": "Internal",
      "sections": [
        {{"title":"Product Vision","body":"[2-3 sentences on the product philosophy, what problem it was built to solve, and the founding insight.]"}},
        {{"title":"Core Capabilities","bullets":["[Capability 1] — [Specific technical feature + user benefit + any metric]","[Capability 2] — [Specific technical feature + user benefit + any metric]","[Capability 3] — [Specific technical feature + user benefit + any metric]","[Capability 4] — [Specific technical feature + user benefit + any metric]","[Capability 5] — [Specific technical feature + user benefit + any metric]"]}},
        {{"title":"Performance Benchmarks","stats":[{{"num":"<100ms","label":"Response Time"}},{{"num":"99.97%","label":"Uptime SLA"}},{{"num":"SOC2","label":"Compliance"}},{{"num":"256-bit","label":"Encryption"}}]}},
        {{"title":"Integration Ecosystem","body":"[2-3 sentences about API, integrations, and how the product fits into existing workflows.]"}},
        {{"title":"Roadmap Highlights","bullets":["Q1: [Upcoming feature and customer benefit]","Q2: [Upcoming feature and customer benefit]","Q3: [Major platform milestone]","Vision: [Long-term product direction]"]}}
      ]
    }},
    {{
      "type": "One-Pager",
      "title": "[Brand] at a Glance",
      "subtitle": "Everything you need to know in under 2 minutes",
      "classification": "Public",
      "sections": [
        {{"title":"What We Do","body":"[2 sentence crystal-clear explanation of who you help, what you help them do, and how.]"}},
        {{"title":"Who We Help","bullets":["[Primary customer profile — industry, size, role]","[Secondary customer profile]","[The specific challenge each faces that you solve]"]}},
        {{"title":"Key Results","stats":[{{"num":"10k+","label":"Customers"}},{{"num":"98%","label":"Retention"}},{{"num":"45m","label":"Onboarding Time"}},{{"num":"4.9★","label":"Rating"}}]}},
        {{"title":"Get Started","body":"[1-2 sentence clear call to action with URL, contact, and any offer or guarantee.]"}}
      ]
    }}
  ]
}}"""


def task_seo(brand_name, headline, subheadline, description) -> str:
    return f"""Generate complete SEO metadata for:
Brand: {brand_name}
Headline: {headline}
Subheadline: {subheadline}
Business: {description}

Return ONLY this JSON:
{{
  "meta_title": "60-65 char SEO title including brand name and primary keyword",
  "meta_description": "150-160 char compelling meta description with CTA and keyword",
  "og_title": "Open Graph title for social sharing, slightly more emotional than meta_title",
  "og_description": "OG description, 120-140 chars, benefit-focused",
  "keywords": ["primary keyword","secondary keyword","keyword 3","keyword 4","keyword 5","keyword 6","keyword 7","keyword 8"],
  "schema_json_ld": {{
    "@context": "https://schema.org",
    "@type": "SoftwareApplication",
    "name": "{brand_name}",
    "description": "Schema description 100-120 chars",
    "applicationCategory": "BusinessApplication",
    "offers": {{"@type":"Offer","price":"0","priceCurrency":"USD"}}
  }},
  "twitter_card": "summary_large_image",
  "canonical_url_path": "/",
  "structured_data_type": "SoftwareApplication",
  "focus_keyword": "1-3 word primary keyword phrase",
  "h1_recommendation": "H1 variant: same intent as headline but slightly different phrasing for SEO"
}}"""


def task_design(description, industry, tone) -> str:
    return f"""Recommend a visual design system for:
Description: {description}
Industry: {industry}
Tone: {tone}

Return ONLY this JSON:
{{
  "brand_color": "#hexcode (primary — vivid, distinctive, not generic blue)",
  "accent_color": "#hexcode (secondary — high contrast against brand_color)",
  "bg_color": "#hexcode (background — should be dark for tech/bold, light for traditional)",
  "text_color": "#hexcode (primary text color — high contrast against bg)",
  "font_pair": "one of: Modern | Editorial | Technical | Bold | Elegant | Futuristic",
  "hero_layout": "one of: Centered | Split | Minimal",
  "animation": "one of: Smooth | Dramatic | None",
  "bg_effect": "one of: Gradient Mesh | Dots | Lines | Solid",
  "design_rationale": "2-3 sentence explanation of why these choices suit the brand",
  "color_palette_name": "Creative name for this palette (2-3 words)",
  "mood_words": ["word1","word2","word3"],
  "competitor_contrast": "1 sentence on how this design differentiates from industry convention"
}}"""


def task_qa(content_json) -> str:
    return f"""Review this landing page content for quality and completeness.
Content: {json.dumps(content_json)[:2000]}...

Return ONLY this JSON:
{{
  "overall_score": 85,
  "headline_score": 90,
  "cta_score": 88,
  "copy_clarity_score": 82,
  "conversion_score": 87,
  "issues": ["Issue 1 if any","Issue 2 if any"],
  "improvements": ["Specific improvement suggestion 1","Specific improvement suggestion 2","Specific improvement suggestion 3"],
  "ab_variants": {{
    "headline_variant_b": "Alternative headline version B",
    "cta_variant_b": "Alternative CTA button text B",
    "subheadline_variant_b": "Alternative subheadline B"
  }},
  "strengths": ["What's working well 1","What's working well 2"],
  "conversion_verdict": "One sentence overall conversion potential assessment"
}}"""


def task_copy(brand_name, headline, description, cta_primary) -> str:
    return f"""Create supplementary marketing copy for:
Brand: {brand_name}. Headline: {headline}. Business: {description}. CTA: {cta_primary}

Return ONLY this JSON:
{{
  "email_subject_lines": [
    "Subject line 1 — curiosity-driven",
    "Subject line 2 — benefit-driven",
    "Subject line 3 — urgency-driven",
    "Subject line 4 — social-proof-driven"
  ],
  "email_sequence": [
    {{"subject":"Welcome email subject","body":"150-word welcome email body. Warm, benefit-focused, sets expectations.","type":"Welcome"}},
    {{"subject":"Day 3 nurture subject","body":"120-word nurture email. Share one key insight or tip. Soft CTA.","type":"Nurture"}},
    {{"subject":"Day 7 conversion subject","body":"130-word conversion email. Address main objection. Strong CTA.","type":"Conversion"}}
  ],
  "ad_copy": [
    {{"platform":"Google Search","headline":"25 char headline","description":"90 char description with CTA","type":"Search Ad"}},
    {{"platform":"Facebook/Instagram","headline":"Attention-grabbing hook","description":"Primary text 80-100 words with story arc","type":"Social Ad"}},
    {{"platform":"LinkedIn","headline":"Professional value prop headline","description":"Business-focused 100-word LinkedIn ad copy","type":"LinkedIn Ad"}}
  ],
  "social_posts": [
    {{"platform":"Twitter/X","text":"Tweet under 240 chars with hook + value + CTA","hashtags":["#tag1","#tag2","#tag3"]}},
    {{"platform":"LinkedIn","text":"150-word LinkedIn post with hook, insight, and CTA","hashtags":["#business","#productivity"]}},
    {{"platform":"Instagram","text":"Caption with hook, story, and CTA under 150 words","hashtags":["#tag1","#tag2","#tag3","#tag4"]}}
  ]
}}"""


# ─────────────────────────────────────────────────────────────────
# SESSION STATE
# ─────────────────────────────────────────────────────────────────

defaults = {
    "page_data": None,
    "pptx_data": None,
    "pdf_docs": None,
    "seo_data": None,
    "design_data": None,
    "qa_data": None,
    "copy_data": None,
    "description": "",
    "tone": "Professional",
    "industry": "Technology",
    "audience": "B2B professionals",
    "brand_color": "#4d6fff",
    "accent_color": "#ff3355",
    "bg_color": "#060810",
    "text_color": "#e4e7f0",
    "font_pair": "Modern",
    "hero_layout": "Centered",
    "animation": "Smooth",
    "bg_effect": "Gradient Mesh",
    "nav_sticky": True,
    "cta_url": "#signup",
    "show_testimonials": True,
    "show_faq": True,
    "show_stats": True,
    "show_cta_banner": True,
    "show_nav": True,
    "gen_count": 0,
    "agent_log": [],
    "agent_statuses": {name: "idle" for name in ["ContentAgent","PptxAgent","PdfAgent","SeoAgent","DesignAgent","QaAgent","CopyAgent"]},
    "orch": None,
    "selected_agents": ["ContentAgent","PptxAgent","PdfAgent","SeoAgent","DesignAgent","QaAgent","CopyAgent"],
}

for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# Build or re-use orchestrator
if st.session_state.orch is None:
    st.session_state.orch = build_orchestrator()

orch: BmAddOrchestrator = st.session_state.orch


# ─────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────

def get_settings():
    return {
        "brand_color": st.session_state.brand_color,
        "accent_color": st.session_state.accent_color,
        "bg_color": st.session_state.bg_color,
        "text_color": st.session_state.text_color,
        "font_pair": st.session_state.font_pair,
        "hero_layout": st.session_state.hero_layout,
        "animation": st.session_state.animation,
        "bg_effect": st.session_state.bg_effect,
        "nav_sticky": st.session_state.nav_sticky,
        "cta_url": st.session_state.cta_url,
        "show_testimonials": st.session_state.show_testimonials,
        "show_faq": st.session_state.show_faq,
        "show_stats": st.session_state.show_stats,
        "show_cta_banner": st.session_state.show_cta_banner,
        "show_nav": st.session_state.show_nav,
    }


def agent_status_badge(name: str) -> str:
    s = st.session_state.agent_statuses.get(name, "idle")
    colors = {"idle": "#3a3f58", "running": "#f59e0b", "done": "#2dd4a0", "error": "#ff4466"}
    labels = {"idle": "IDLE", "running": "RUNNING", "done": "DONE", "error": "ERROR"}
    return f'<span style="background:{colors[s]}22;border:1px solid {colors[s]}55;color:{colors[s]};font-family:\'Space Mono\',monospace;font-size:0.6rem;padding:2px 7px;border-radius:3px;letter-spacing:0.08em">{labels[s]}</span>'


# ─────────────────────────────────────────────────────────────────
# ════════════════════════════════════════════════════════════════
# SIDEBAR
# ════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown("""
    <div class="pf-header">
        <div class="pf-logo">PAGE<span>FORGE</span><sub>PRO</sub></div>
        <div class="pf-tagline">AI Document Studio · BmAdd Agents</div>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("<div style='padding: 1rem 1.1rem 0'>", unsafe_allow_html=True)
    
    # API KEY
    st.markdown('<span class="section-label">🔑 OpenAI API Key</span>', unsafe_allow_html=True)
    api_key = st.text_input("API Key", type="password", placeholder="sk-...", label_visibility="collapsed")
    
    model_choice = st.selectbox("Model", ["gpt-4o-mini", "gpt-4o", "gpt-3.5-turbo"], label_visibility="visible")
    
    st.markdown("<hr>", unsafe_allow_html=True)
    
    # GENERATION STATS
    pdata = st.session_state.page_data
    pptx_d = st.session_state.pptx_data
    pdf_d = st.session_state.pdf_docs
    
    st.markdown(f"""
    <div class="stat-grid-3">
        <div class="stat-box"><div class="n">{st.session_state.gen_count}</div><div class="l">Runs</div></div>
        <div class="stat-box"><div class="n">{'✓' if pdata else '—'}</div><div class="l">HTML</div></div>
        <div class="stat-box"><div class="n">{len(pptx_d.get('slides',[])) if isinstance(pptx_d,dict) else '—'}</div><div class="l">Slides</div></div>
    </div>
    <div class="stat-grid-3">
        <div class="stat-box"><div class="n">{len(pdf_d) if isinstance(pdf_d,list) else '—'}</div><div class="l">PDFs</div></div>
        <div class="stat-box"><div class="n">{'✓' if st.session_state.seo_data else '—'}</div><div class="l">SEO</div></div>
        <div class="stat-box"><div class="n">{'✓' if st.session_state.qa_data else '—'}</div><div class="l">QA</div></div>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("<hr>", unsafe_allow_html=True)
    
    # AGENT REGISTRY
    st.markdown('<span class="section-label">🤖 BmAdd Agents</span>', unsafe_allow_html=True)
    
    agent_info = [
        ("ContentAgent", "✍️", "Copy & headlines"),
        ("PptxAgent", "📊", "Slide deck"),
        ("PdfAgent", "📄", "PDF suite"),
        ("SeoAgent", "🔍", "SEO metadata"),
        ("DesignAgent", "🎨", "Design system"),
        ("QaAgent", "✅", "QA & scoring"),
        ("CopyAgent", "📢", "Email & ads"),
    ]
    
    for name, emoji, role_short in agent_info:
        status = st.session_state.agent_statuses.get(name, "idle")
        s_colors = {"idle": "#3a3f58", "running": "#f59e0b", "done": "#2dd4a0", "error": "#ff4466"}
        s_color = s_colors.get(status, "#3a3f58")
        st.markdown(f"""
        <div class="agent-card {status}">
            <div class="agent-name">
                <span>{emoji} {name}</span>
                <span style="font-family:'Space Mono',monospace;font-size:0.58rem;color:{s_color};letter-spacing:0.08em">{status.upper()}</span>
            </div>
            <div class="agent-role">{role_short}</div>
        </div>""", unsafe_allow_html=True)
    
    st.markdown("<hr>", unsafe_allow_html=True)
    
    # QUICK GUIDE
    st.markdown('<span class="section-label">📋 How to use</span>', unsafe_allow_html=True)
    steps = [
        ("1", "Enter business description + API key"),
        ("2", "Select agents to run"),
        ("3", "Click Launch Agent Swarm"),
        ("4", "Edit content in Copy / Design tabs"),
        ("5", "Download HTML, PPTX, PDFs & more"),
    ]
    for num, text in steps:
        st.markdown(f"""<div class="step-block">
            <div class="step-num">{num}</div>
            <div class="step-content"><div class="step-desc">{text}</div></div>
        </div>""", unsafe_allow_html=True)
    
    st.markdown("</div>", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────
# ════════════════════════════════════════════════════════════════
# MAIN AREA
# ════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────

st.markdown("""
<div class="page-title">PAGE<span class="accent">FORGE</span> PRO</div>
<p style="font-family:'Space Mono',monospace;font-size:0.68rem;color:#3a3f58;letter-spacing:0.18em;text-transform:uppercase;margin-bottom:1.5rem">
AI Document Studio — HTML · PPTX · PDF · SEO · BmAdd Agents
</p>
""", unsafe_allow_html=True)


# ─── INPUT + CONFIG ───
with st.expander("⚙️  Business Description & Configuration", expanded=True):
    col_desc, col_meta, col_agents = st.columns([2.2, 1, 1])
    
    with col_desc:
        st.markdown('<span class="section-label">💡 Describe Your Business</span>', unsafe_allow_html=True)
        description = st.text_area(
            "desc", value=st.session_state.description, height=130,
            placeholder="Describe your product, service, or business in detail. Include your unique value proposition, target market, key benefits, competitive advantages, and any specific messaging requirements. The more detail you provide, the better the output quality...",
            label_visibility="collapsed"
        )
        st.session_state.description = description
        
        col_t, col_i, col_a = st.columns(3)
        with col_t:
            st.session_state.tone = st.selectbox("Tone", ["Professional", "Playful", "Bold & Edgy", "Minimalist", "Luxury", "Technical", "Friendly", "Urgent"])
        with col_i:
            st.session_state.industry = st.selectbox("Industry", ["Technology", "SaaS", "E-commerce", "Healthcare", "Finance", "Education", "Creative Agency", "Real Estate", "Food & Bev", "Consulting", "Other"])
        with col_a:
            st.session_state.audience = st.text_input("Target Audience", value=st.session_state.audience, placeholder="e.g. Marketing VPs at B2B SaaS")
    
    with col_meta:
        st.markdown('<span class="section-label">🎨 Quick Colors</span>', unsafe_allow_html=True)
        st.session_state.brand_color = st.color_picker("Brand Color", st.session_state.brand_color)
        st.session_state.accent_color = st.color_picker("Accent Color", st.session_state.accent_color)
        st.session_state.bg_color = st.color_picker("Background", st.session_state.bg_color)
        st.session_state.text_color = st.color_picker("Text Color", st.session_state.text_color)
        
        st.markdown('<span class="section-label">🔗 CTA URL</span>', unsafe_allow_html=True)
        st.session_state.cta_url = st.text_input("CTA URL", value=st.session_state.cta_url, label_visibility="collapsed")
    
    with col_agents:
        st.markdown('<span class="section-label">🤖 Select Agents</span>', unsafe_allow_html=True)
        all_agents = ["ContentAgent","PptxAgent","PdfAgent","SeoAgent","DesignAgent","QaAgent","CopyAgent"]
        agent_emojis = {"ContentAgent":"✍️","PptxAgent":"📊","PdfAgent":"📄","SeoAgent":"🔍","DesignAgent":"🎨","QaAgent":"✅","CopyAgent":"📢"}
        selected = []
        for ag in all_agents:
            checked = st.checkbox(f"{agent_emojis[ag]} {ag}", value=(ag in st.session_state.selected_agents), key=f"agent_chk_{ag}")
            if checked:
                selected.append(ag)
        st.session_state.selected_agents = selected


# ─── LAUNCH BUTTON ───
col_launch, col_status = st.columns([1, 3])
with col_launch:
    launch = st.button("⚡ Launch BmAdd Agent Swarm", use_container_width=True)
with col_status:
    if not api_key:
        st.info("🔑 Add your OpenAI API key in the sidebar to run the agent swarm")
    elif not st.session_state.description.strip():
        st.info("📝 Enter a business description above to get started")
    elif not st.session_state.selected_agents:
        st.warning("🤖 Select at least one agent to run")
    else:
        st.success(f"✅ Ready — {len(st.session_state.selected_agents)} agents selected · {model_choice} · Click Launch to begin")


# ─────────────────────────────────────────────────────────────────
# GENERATION PIPELINE
# ─────────────────────────────────────────────────────────────────

if launch:
    if not api_key:
        st.error("❌ OpenAI API key required. Add it in the sidebar.")
    elif not st.session_state.description.strip():
        st.error("❌ Please enter a business description.")
    elif not st.session_state.selected_agents:
        st.error("❌ Select at least one agent.")
    else:
        # Reset
        orch.reset()
        st.session_state.agent_log = []
        st.session_state.agent_statuses = {n: "idle" for n in defaults["agent_statuses"]}
        
        # Override model
        for agent in orch.agents.values():
            # Patch model in run_sync via monkey-patching
            pass
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        log_container = st.empty()
        
        total_agents = len(st.session_state.selected_agents)
        completed = 0
        
        def update_log_display():
            log_html = '<div style="background:var(--surface2,#111320);border:1px solid var(--border,#1a1e32);border-radius:8px;padding:10px;max-height:220px;overflow-y:auto;font-family:\'Space Mono\',monospace;font-size:0.68rem;">'
            for entry in st.session_state.agent_log[-15:]:
                s = entry.get("status","info")
                c = {"running":"#f59e0b","done":"#2dd4a0","error":"#ff4466","info":"#4d6fff"}.get(s,"#525870")
                log_html += f'<div style="padding:3px 0;border-bottom:1px solid rgba(255,255,255,0.04);"><span style="color:#3a3f58">{entry.get("ts","")} </span><span style="color:{c};font-weight:700">[{entry.get("agent","")}]</span> <span style="color:#e4e7f0">{entry.get("msg","")}</span></div>'
            log_html += "</div>"
            log_container.markdown(log_html, unsafe_allow_html=True)
        
        # ── AGENT 1: ContentAgent ──
        if "ContentAgent" in st.session_state.selected_agents:
            status_text.markdown("⏳ **ContentAgent** — Generating landing page copy...")
            task = task_content(description, st.session_state.tone, st.session_state.industry, st.session_state.audience)
            # Inject model into task by running agent's run_sync directly
            agent = orch.agents["ContentAgent"]
            
            import urllib.request as _ureq
            import json as _json
            
            # Patch model
            original_run = agent.run_sync
            
            def patched_run(task, key, callback=None, _agent=agent, _model=model_choice):
                _agent.status = "running"
                _agent.start_time = time.time()
                _agent.memory.write("task", task)
                if callback:
                    callback(_agent.name, "running", f"Starting: {task[:55]}...")
                try:
                    payload = _json.dumps({
                        "model": _model,
                        "messages": [
                            {"role": "system", "content": _agent._build_system_prompt()},
                            {"role": "user", "content": task}
                        ],
                        "temperature": 0.8,
                        "max_tokens": 3500
                    }).encode("utf-8")
                    req = _ureq.Request(
                        "https://api.openai.com/v1/chat/completions",
                        data=payload,
                        headers={"Content-Type": "application/json", "Authorization": f"Bearer {key}"},
                        method="POST"
                    )
                    with _ureq.urlopen(req, timeout=55) as resp:
                        result = _json.loads(resp.read().decode("utf-8"))
                        raw = result["choices"][0]["message"]["content"]
                        _agent.token_count = result.get("usage", {}).get("total_tokens", 0)
                    raw_clean = raw.strip()
                    raw_clean = re.sub(r'^```(?:json)?\s*', '', raw_clean, flags=re.MULTILINE)
                    raw_clean = re.sub(r'\s*```\s*$', '', raw_clean, flags=re.MULTILINE).strip()
                    json_match = re.search(r'\{.*\}', raw_clean, re.DOTALL)
                    if json_match:
                        _agent.output = _json.loads(json_match.group())
                    else:
                        _agent.output = _json.loads(raw_clean)
                    _agent.status = "done"
                    _agent.end_time = time.time()
                    if callback:
                        callback(_agent.name, "done", f"Completed in {_agent.elapsed()} · {_agent.token_count} tokens")
                    return _agent.output
                except Exception as e:
                    _agent.error_msg = str(e)
                    _agent.status = "error"
                    _agent.end_time = time.time()
                    if callback:
                        callback(_agent.name, "error", f"Error: {str(e)[:120]}")
                    return None
            
            # Use the patched run for all agents
            def run_agent_with_model(agent_name, task, api_key, model):
                agent = orch.agents[agent_name]
                agent.status = "running"
                agent.start_time = time.time()
                agent.memory.write("task", task)
                
                if "agent_statuses" in st.session_state:
                    st.session_state.agent_statuses[agent_name] = "running"
                
                callback = orch._add_log
                if callback:
                    callback(agent.name, "running", f"Processing...")
                
                try:
                    payload = _json.dumps({
                        "model": model,
                        "messages": [
                            {"role": "system", "content": agent._build_system_prompt()},
                            {"role": "user", "content": task}
                        ],
                        "temperature": 0.8,
                        "max_tokens": 3500
                    }).encode("utf-8")
                    req = _ureq.Request(
                        "https://api.openai.com/v1/chat/completions",
                        data=payload,
                        headers={"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"},
                        method="POST"
                    )
                    with _ureq.urlopen(req, timeout=60) as resp:
                        result = _json.loads(resp.read().decode("utf-8"))
                        raw = result["choices"][0]["message"]["content"]
                        agent.token_count = result.get("usage", {}).get("total_tokens", 0)
                    
                    raw_clean = raw.strip()
                    raw_clean = re.sub(r'^```(?:json)?\s*', '', raw_clean, flags=re.MULTILINE)
                    raw_clean = re.sub(r'\s*```\s*$', '', raw_clean, flags=re.MULTILINE).strip()
                    
                    # Try to find JSON
                    json_match = re.search(r'\{.*\}', raw_clean, re.DOTALL)
                    if json_match:
                        agent.output = _json.loads(json_match.group())
                    else:
                        agent.output = _json.loads(raw_clean)
                    
                    orch.context[agent_name] = agent.output
                    orch.total_tokens += agent.token_count
                    agent.status = "done"
                    agent.end_time = time.time()
                    st.session_state.agent_statuses[agent_name] = "done"
                    
                    if callback:
                        callback(agent.name, "done", f"✓ Done in {agent.elapsed()} ({agent.token_count} tokens)")
                    
                    return agent.output
                    
                except _ureq.HTTPError as e:
                    err = e.read().decode("utf-8")[:200]
                    agent.error_msg = err
                    agent.status = "error"
                    agent.end_time = time.time()
                    st.session_state.agent_statuses[agent_name] = "error"
                    if callback:
                        callback(agent.name, "error", f"HTTP {e.code}: {err[:100]}")
                    return None
                except Exception as e:
                    agent.error_msg = str(e)
                    agent.status = "error"
                    agent.end_time = time.time()
                    st.session_state.agent_statuses[agent_name] = "error"
                    if callback:
                        callback(agent.name, "error", f"Error: {str(e)[:120]}")
                    return None
            
            # Run ContentAgent
            result = run_agent_with_model("ContentAgent", task, api_key, model_choice)
            if result:
                st.session_state.page_data = result
                completed += 1
                progress_bar.progress(int(completed / total_agents * 100))
                update_log_display()
        
        # ── AGENT 2: DesignAgent ──
        if "DesignAgent" in st.session_state.selected_agents:
            status_text.markdown("⏳ **DesignAgent** — Crafting design system...")
            task = task_design(description, st.session_state.industry, st.session_state.tone)
            result = run_agent_with_model("DesignAgent", task, api_key, model_choice)
            if result:
                st.session_state.design_data = result
                # Apply design recommendations
                if result.get("brand_color"):
                    st.session_state.brand_color = result["brand_color"]
                if result.get("accent_color"):
                    st.session_state.accent_color = result["accent_color"]
                if result.get("bg_color"):
                    st.session_state.bg_color = result["bg_color"]
                if result.get("text_color"):
                    st.session_state.text_color = result["text_color"]
                if result.get("font_pair"):
                    st.session_state.font_pair = result["font_pair"]
                if result.get("hero_layout"):
                    st.session_state.hero_layout = result["hero_layout"]
                if result.get("animation"):
                    st.session_state.animation = result["animation"]
                if result.get("bg_effect"):
                    st.session_state.bg_effect = result["bg_effect"]
            completed += 1
            progress_bar.progress(int(completed / total_agents * 100))
            update_log_display()
        
        # ── AGENT 3: SeoAgent ──
        if "SeoAgent" in st.session_state.selected_agents and st.session_state.page_data:
            status_text.markdown("⏳ **SeoAgent** — Generating SEO metadata...")
            d = st.session_state.page_data
            task = task_seo(d.get("brand_name","Brand"), d.get("headline",""), d.get("subheadline",""), description)
            result = run_agent_with_model("SeoAgent", task, api_key, model_choice)
            if result:
                st.session_state.seo_data = result
                # Merge SEO into page data
                if st.session_state.page_data:
                    st.session_state.page_data["seo_meta"] = result
            completed += 1
            progress_bar.progress(int(completed / total_agents * 100))
            update_log_display()
        
        # ── AGENT 4: PptxAgent ──
        if "PptxAgent" in st.session_state.selected_agents and st.session_state.page_data:
            status_text.markdown("⏳ **PptxAgent** — Architecting slide deck...")
            d = st.session_state.page_data
            task = task_pptx(description, d.get("brand_name","Brand"), d.get("headline",""), d.get("features",[]))
            result = run_agent_with_model("PptxAgent", task, api_key, model_choice)
            if result:
                st.session_state.pptx_data = result
            completed += 1
            progress_bar.progress(int(completed / total_agents * 100))
            update_log_display()
        
        # ── AGENT 5: PdfAgent ──
        if "PdfAgent" in st.session_state.selected_agents and st.session_state.page_data:
            status_text.markdown("⏳ **PdfAgent** — Creating PDF document suite...")
            d = st.session_state.page_data
            task = task_pdfs(description, d.get("brand_name","Brand"), d.get("headline",""),
                            d.get("features",[]), d.get("testimonials",[]), d.get("stats",[]))
            result = run_agent_with_model("PdfAgent", task, api_key, model_choice)
            if result and result.get("documents"):
                st.session_state.pdf_docs = result["documents"]
            completed += 1
            progress_bar.progress(int(completed / total_agents * 100))
            update_log_display()
        
        # ── AGENT 6: QaAgent ──
        if "QaAgent" in st.session_state.selected_agents and st.session_state.page_data:
            status_text.markdown("⏳ **QaAgent** — Quality reviewing all content...")
            task = task_qa(st.session_state.page_data)
            result = run_agent_with_model("QaAgent", task, api_key, model_choice)
            if result:
                st.session_state.qa_data = result
            completed += 1
            progress_bar.progress(int(completed / total_agents * 100))
            update_log_display()
        
        # ── AGENT 7: CopyAgent ──
        if "CopyAgent" in st.session_state.selected_agents and st.session_state.page_data:
            status_text.markdown("⏳ **CopyAgent** — Writing email sequences & ad copy...")
            d = st.session_state.page_data
            task = task_copy(d.get("brand_name","Brand"), d.get("headline",""), description, d.get("cta_primary",""))
            result = run_agent_with_model("CopyAgent", task, api_key, model_choice)
            if result:
                st.session_state.copy_data = result
            completed += 1
            progress_bar.progress(int(completed / total_agents * 100))
            update_log_display()
        
        progress_bar.progress(100)
        st.session_state.gen_count += 1
        
        orch._add_log("BmAddOrchestrator", "done", 
                     f"Pipeline complete — {completed}/{total_agents} agents succeeded · {orch.total_tokens} total tokens")
        update_log_display()
        
        status_text.success(f"✅ Done! {completed}/{total_agents} agents completed · {orch.total_tokens} tokens used")
        time.sleep(0.5)
        st.rerun()


# ─────────────────────────────────────────────────────────────────
# EDITOR & OUTPUT TABS (only when data exists)
# ─────────────────────────────────────────────────────────────────

if st.session_state.page_data:
    data = st.session_state.page_data
    settings = get_settings()
    
    st.markdown("<hr>", unsafe_allow_html=True)
    
    tabs = st.tabs([
        "✏️ Copy Editor", "🎨 Design & Style", "📊 Slide Deck", 
        "📄 PDF Suite", "🔍 SEO", "✅ QA Report", 
        "📢 Email & Ads", "💻 Code", "👁️ Preview", "📦 Export"
    ])
    
    # ════════════════════════════════════════════════════════════════
    # TAB 1 — COPY EDITOR
    # ════════════════════════════════════════════════════════════════
    with tabs[0]:
        st.markdown('<div class="chip chip-red">✏️ COPY EDITOR</div>', unsafe_allow_html=True)
        
        # Hero content
        col_h1, col_h2 = st.columns(2)
        with col_h1:
            data['headline'] = st.text_input("🔥 Main Headline", value=data.get('headline',''))
            data['subheadline'] = st.text_input("📌 Subheadline", value=data.get('subheadline',''))
            data['brand_name'] = st.text_input("🏷️ Brand Name", value=data.get('brand_name',''))
            data['tagline'] = st.text_input("✨ Tagline", value=data.get('tagline',''))
            data['social_proof'] = st.text_input("🏆 Social Proof Badge", value=data.get('social_proof',''))
        with col_h2:
            data['cta_primary'] = st.text_input("🔘 Primary CTA Button", value=data.get('cta_primary','Get Started'))
            data['cta_secondary'] = st.text_input("🔲 Secondary CTA Button", value=data.get('cta_secondary','Learn More'))
            st.session_state.cta_url = st.text_input("🔗 CTA URL", value=st.session_state.cta_url)
            data['email_placeholder'] = st.text_input("📧 Email Placeholder", value=data.get('email_placeholder',''))
            data['footer_text'] = st.text_input("👇 Footer Tagline", value=data.get('footer_text',''))
        
        data['hero_body'] = st.text_area("📝 Hero Body Text", value=data.get('hero_body',''), height=100)
        
        # QA SUGGESTED VARIANTS
        qa = st.session_state.qa_data
        if qa and qa.get("ab_variants"):
            with st.expander("🧪 A/B Variants from QaAgent", expanded=False):
                ab = qa["ab_variants"]
                c1, c2 = st.columns(2)
                with c1:
                    if ab.get("headline_variant_b"):
                        st.markdown(f"**Headline B:** {ab['headline_variant_b']}")
                        if st.button("Use Headline B", key="use_hb"):
                            data["headline"] = ab["headline_variant_b"]
                            st.rerun()
                    if ab.get("cta_variant_b"):
                        st.markdown(f"**CTA B:** {ab['cta_variant_b']}")
                        if st.button("Use CTA B", key="use_ctab"):
                            data["cta_primary"] = ab["cta_variant_b"]
                            st.rerun()
                with c2:
                    if ab.get("subheadline_variant_b"):
                        st.markdown(f"**Subheadline B:** {ab['subheadline_variant_b']}")
                        if st.button("Use Subheadline B", key="use_shb"):
                            data["subheadline"] = ab["subheadline_variant_b"]
                            st.rerun()
        
        st.markdown("<hr>", unsafe_allow_html=True)
        st.markdown('<div class="chip chip-blue">⚡ FEATURES</div>', unsafe_allow_html=True)
        
        features = data.get('features', [])
        cols_feat = st.columns(2)
        new_features = []
        for i, feat in enumerate(features):
            col = cols_feat[i % 2]
            with col:
                with st.expander(f"{feat.get('icon','⚡')} {feat.get('title', f'Feature {i+1}')}", expanded=False):
                    c1, c2 = st.columns([1, 4])
                    with c1:
                        icon = st.text_input("Icon", value=feat.get('icon','⚡'), key=f"feat_icon_{i}")
                    with c2:
                        title = st.text_input("Title", value=feat.get('title',''), key=f"feat_title_{i}")
                    desc = st.text_area("Description", value=feat.get('desc',''), key=f"feat_desc_{i}", height=80)
                    new_features.append({"icon": icon, "title": title, "desc": desc})
        data['features'] = new_features
        
        col_add, col_rem = st.columns(2)
        with col_add:
            if st.button("➕ Add Feature"):
                data['features'].append({"icon": "⭐", "title": "New Feature", "desc": "Describe this feature."})
                st.rerun()
        
        st.markdown("<hr>", unsafe_allow_html=True)
        st.markdown('<div class="chip chip-teal">💬 TESTIMONIALS</div>', unsafe_allow_html=True)
        
        testimonials = data.get('testimonials', [])
        for i, t in enumerate(testimonials):
            with st.expander(f"💬 {t.get('name', f'Testimonial {i+1}')}", expanded=False):
                c1, c2, c3 = st.columns([2, 2, 1])
                with c1:
                    t['name'] = st.text_input("Name", value=t.get('name',''), key=f"test_name_{i}")
                with c2:
                    t['role'] = st.text_input("Role / Company", value=t.get('role',''), key=f"test_role_{i}")
                with c3:
                    t['rating'] = st.slider("⭐", 1, 5, t.get('rating',5), key=f"test_rating_{i}")
                t['text'] = st.text_area("Testimonial Text", value=t.get('text',''), key=f"test_text_{i}", height=80)
        
        if st.button("➕ Add Testimonial"):
            data['testimonials'].append({"name":"New Person","role":"Title, Company","text":"Amazing product.","rating":5})
            st.rerun()
        
        st.markdown("<hr>", unsafe_allow_html=True)
        st.markdown('<div class="chip chip-gold">❓ FAQ</div>', unsafe_allow_html=True)
        
        faq = data.get('faq', [])
        for i, item in enumerate(faq):
            with st.expander(f"❓ {item.get('q','FAQ')[:55]}...", expanded=False):
                item['q'] = st.text_input("Question", value=item.get('q',''), key=f"faq_q_{i}")
                item['a'] = st.text_area("Answer", value=item.get('a',''), key=f"faq_a_{i}", height=80)
        
        if st.button("➕ Add FAQ"):
            data['faq'].append({"q":"New question?","a":"Answer here."})
            st.rerun()
        
        st.markdown("<hr>", unsafe_allow_html=True)
        st.markdown('<div class="chip chip-purple">📊 STATS</div>', unsafe_allow_html=True)
        
        stats = data.get('stats', [])
        stat_cols = st.columns(4)
        for i, stat in enumerate(stats):
            with stat_cols[i % 4]:
                stat['num'] = st.text_input("Number", value=stat.get('num',''), key=f"stat_num_{i}")
                stat['label'] = st.text_input("Label", value=stat.get('label',''), key=f"stat_label_{i}")
        
        st.markdown("<hr>", unsafe_allow_html=True)
        st.markdown('<div class="chip chip-blue">🧭 NAVIGATION</div>', unsafe_allow_html=True)
        nav_str = ", ".join(data.get('nav_links', ['Features','Testimonials','Pricing','FAQ']))
        nav_edit = st.text_input("Nav Links (comma-separated)", value=nav_str)
        data['nav_links'] = [l.strip() for l in nav_edit.split(',') if l.strip()]
        
        st.session_state.page_data = data
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 2 — DESIGN & STYLE
    # ════════════════════════════════════════════════════════════════
    with tabs[1]:
        st.markdown('<div class="chip chip-blue">🎨 DESIGN SYSTEM</div>', unsafe_allow_html=True)
        
        # Show DesignAgent recommendation
        dd = st.session_state.design_data
        if dd:
            with st.expander(f"🎨 DesignAgent Recommendation — {dd.get('color_palette_name','Custom Palette')}", expanded=True):
                c1, c2 = st.columns([2, 1])
                with c1:
                    st.markdown(f"**Rationale:** {dd.get('design_rationale','')}")
                    if dd.get('mood_words'):
                        st.markdown(f"**Mood:** {' · '.join(dd['mood_words'])}")
                    if dd.get('competitor_contrast'):
                        st.markdown(f"**Differentiation:** {dd['competitor_contrast']}")
                with c2:
                    st.markdown(f"**Font:** {dd.get('font_pair','')}")
                    st.markdown(f"**Layout:** {dd.get('hero_layout','')}")
                    st.markdown(f"**Animation:** {dd.get('animation','')}")
        
        col_d1, col_d2, col_d3 = st.columns(3)
        
        with col_d1:
            st.markdown("**🎨 Colors**")
            st.session_state.brand_color = st.color_picker("Brand Color", st.session_state.brand_color, key="bc2")
            st.session_state.accent_color = st.color_picker("Accent Color", st.session_state.accent_color, key="ac2")
            st.session_state.bg_color = st.color_picker("Background", st.session_state.bg_color, key="bgc2")
            st.session_state.text_color = st.color_picker("Text Color", st.session_state.text_color, key="tc2")
            
            st.markdown("**Quick Color Themes**")
            themes = {
                "Midnight Blue": ("#4d6fff","#ff3355","#060810","#e4e7f0"),
                "Cyber Green": ("#00d88a","#ff6600","#040d08","#d4f7e4"),
                "Sunset Coral": ("#ff5e5b","#ffb400","#1a0a08","#fff0ee"),
                "Deep Purple": ("#9333ea","#ec4899","#0a0514","#f0e6ff"),
                "Ocean Teal": ("#0891b2","#f59e0b","#040e14","#e0f7ff"),
                "Minimal Light": ("#1a1a2e","#ff3355","#fafafa","#1a1a2e"),
                "Luxury Gold": ("#c9953d","#ff3355","#0a0804","#f5efe0"),
                "Forest Sage": ("#2d9b5e","#ff7043","#080f0c","#e0f2e9"),
            }
            tcols = st.columns(4)
            for i, (tname, tcolors) in enumerate(themes.items()):
                with tcols[i % 4]:
                    st.markdown(f"""<div class="theme-swatch" title="{tname}" onclick="void(0)">
                        <div style="background:{tcolors[0]}"></div>
                        <div style="background:{tcolors[1]}"></div>
                        <div style="background:{tcolors[2]}"></div>
                    </div>""", unsafe_allow_html=True)
                    if st.button(tname[:8], key=f"theme_btn_{tname}"):
                        st.session_state.brand_color = tcolors[0]
                        st.session_state.accent_color = tcolors[1]
                        st.session_state.bg_color = tcolors[2]
                        st.session_state.text_color = tcolors[3]
                        st.rerun()
        
        with col_d2:
            st.markdown("**🔤 Typography**")
            fp_opts = ["Modern", "Editorial", "Technical", "Bold", "Elegant", "Futuristic"]
            st.session_state.font_pair = st.radio("Font Pair", fp_opts,
                index=fp_opts.index(st.session_state.font_pair) if st.session_state.font_pair in fp_opts else 0)
            
            st.markdown("**📐 Hero Layout**")
            hl_opts = ["Centered", "Split", "Minimal"]
            st.session_state.hero_layout = st.radio("Hero Layout", hl_opts,
                index=hl_opts.index(st.session_state.hero_layout) if st.session_state.hero_layout in hl_opts else 0)
        
        with col_d3:
            st.markdown("**✨ Animation**")
            an_opts = ["Smooth", "Dramatic", "None"]
            st.session_state.animation = st.radio("Animation", an_opts,
                index=an_opts.index(st.session_state.animation) if st.session_state.animation in an_opts else 0)
            
            st.markdown("**🌌 Background**")
            bg_opts = ["Gradient Mesh", "Dots", "Lines", "Solid"]
            st.session_state.bg_effect = st.radio("Background Effect", bg_opts,
                index=bg_opts.index(st.session_state.bg_effect) if st.session_state.bg_effect in bg_opts else 0)
            
            st.markdown("**🧩 Sections**")
            st.session_state.show_nav = st.checkbox("Nav Bar", value=st.session_state.show_nav)
            st.session_state.nav_sticky = st.checkbox("Sticky Nav", value=st.session_state.nav_sticky)
            st.session_state.show_stats = st.checkbox("Stats Section", value=st.session_state.show_stats)
            st.session_state.show_testimonials = st.checkbox("Testimonials", value=st.session_state.show_testimonials)
            st.session_state.show_faq = st.checkbox("FAQ", value=st.session_state.show_faq)
            st.session_state.show_cta_banner = st.checkbox("CTA Banner", value=st.session_state.show_cta_banner)
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 3 — SLIDE DECK
    # ════════════════════════════════════════════════════════════════
    with tabs[2]:
        st.markdown('<div class="chip chip-blue">📊 SLIDE DECK EDITOR</div>', unsafe_allow_html=True)
        
        pptx_d = st.session_state.pptx_data
        if not pptx_d:
            st.info("🤖 Run PptxAgent to generate slide deck content first.")
        else:
            slides = pptx_d.get("slides", [])
            
            col_pp, col_pi = st.columns([1, 2])
            with col_pp:
                st.markdown(f"**{len(slides)} slides generated**")
                st.markdown('<div class="tip-block">Slides are exported as a real .pptx file using python-pptx with branded layouts, speaker notes, and theme colors.</div>', unsafe_allow_html=True)
            with col_pi:
                pptx_font = st.selectbox("PPTX Heading Font", ["Calibri", "Georgia", "Arial Black", "Cambria", "Palatino", "Trebuchet MS"], key="pptx_font")
            
            st.markdown("**Edit Slides**")
            updated_slides = []
            for i, slide in enumerate(slides):
                with st.expander(f"Slide {i+1}: {slide.get('title','')}", expanded=False):
                    c1, c2 = st.columns([3, 1])
                    with c1:
                        slide['title'] = st.text_input("Title", value=slide.get('title',''), key=f"slide_title_{i}")
                    with c2:
                        layout_opts = ["content", "stats", "cover", "closing"]
                        cur_layout = slide.get('layout','content')
                        if cur_layout not in layout_opts:
                            cur_layout = "content"
                        slide['layout'] = st.selectbox("Layout", layout_opts, 
                            index=layout_opts.index(cur_layout), key=f"slide_layout_{i}")
                    
                    if slide.get('bullets'):
                        bullets_str = "\n".join(slide['bullets'])
                        new_bullets = st.text_area("Bullets (one per line)", value=bullets_str, 
                                                   key=f"slide_bullets_{i}", height=100)
                        slide['bullets'] = [b.strip() for b in new_bullets.split('\n') if b.strip()]
                    elif slide.get('body'):
                        slide['body'] = st.text_area("Body", value=slide.get('body',''), 
                                                      key=f"slide_body_{i}", height=100)
                    elif slide.get('cta'):
                        slide['cta'] = st.text_area("CTA Text", value=slide.get('cta',''), 
                                                     key=f"slide_cta_{i}", height=80)
                    elif slide.get('subtitle'):
                        slide['subtitle'] = st.text_area("Subtitle", value=slide.get('subtitle',''), 
                                                          key=f"slide_subtitle_{i}", height=60)
                    
                    if slide.get('stats'):
                        for si, stat in enumerate(slide['stats']):
                            sc1, sc2 = st.columns(2)
                            with sc1:
                                stat['num'] = st.text_input("Stat #", value=stat.get('num',''), key=f"slide_statnum_{i}_{si}")
                            with sc2:
                                stat['label'] = st.text_input("Label", value=stat.get('label',''), key=f"slide_statlabel_{i}_{si}")
                    
                    if slide.get('speaker_notes'):
                        slide['speaker_notes'] = st.text_area("Speaker Notes", 
                            value=slide.get('speaker_notes',''), key=f"slide_notes_{i}", height=70)
                    
                    if slide.get('icon') is not None:
                        slide['icon'] = st.text_input("Slide Icon", value=slide.get('icon',''), key=f"slide_icon_{i}")
                    
                    updated_slides.append(slide)
            
            st.session_state.pptx_data['slides'] = updated_slides
            
            st.markdown("<hr>", unsafe_allow_html=True)
            if st.button("➕ Add Slide"):
                st.session_state.pptx_data['slides'].append({
                    "title": "New Slide", "body": "Content here.", "layout": "content", "speaker_notes": ""
                })
                st.rerun()
            
            # SLIDE PREVIEWS
            st.markdown('<span class="section-label">Slide Previews</span>', unsafe_allow_html=True)
            prev_cols = st.columns(5)
            for i, slide in enumerate(slides[:10]):
                with prev_cols[i % 5]:
                    is_cover = i == 0
                    bg = st.session_state.brand_color if is_cover else "#13151f"
                    title_color = "#ffffff" if is_cover else "#e4e7f0"
                    bar_color = st.session_state.accent_color if is_cover else st.session_state.brand_color
                    bullets_preview = ""
                    for b in (slide.get('bullets') or [])[:3]:
                        bullets_preview += f'<div style="font-size:0.52rem;color:{"rgba(255,255,255,.65)" if is_cover else "#525870"};margin-bottom:2px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap">• {b}</div>'
                    if slide.get('body'):
                        bullets_preview = f'<div style="font-size:0.52rem;color:{"rgba(255,255,255,.65)" if is_cover else "#525870"};line-height:1.4">{slide["body"][:80]}...</div>'
                    st.markdown(f"""
                    <div class="slide-preview" style="background:{bg}">
                        <div class="slide-preview-num">{i+1}</div>
                        <div class="slide-preview-title" style="color:{title_color}">{(slide.get('title',''))[:35]}</div>
                        {bullets_preview}
                        <div class="slide-preview-bar" style="background:{bar_color}"></div>
                    </div>""", unsafe_allow_html=True)
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 4 — PDF SUITE
    # ════════════════════════════════════════════════════════════════
    with tabs[3]:
        st.markdown('<div class="chip chip-orange">📄 PDF DOCUMENT SUITE</div>', unsafe_allow_html=True)
        
        pdf_d = st.session_state.pdf_docs
        if not pdf_d:
            st.info("🤖 Run PdfAgent to generate the PDF document suite first.")
        else:
            st.markdown(f"**{len(pdf_d)} documents generated by PdfAgent**")
            st.markdown('<div class="tip-block">Each document is rendered as a professional PDF using reportlab with brand colors, typography, tables, stat callouts, and proper page layout.</div>', unsafe_allow_html=True)
            
            for di, doc in enumerate(pdf_d):
                with st.expander(f"📄 {doc.get('type','Document')} — {doc.get('title','')}", expanded=False):
                    c1, c2 = st.columns([2, 1])
                    with c1:
                        doc['title'] = st.text_input("Document Title", value=doc.get('title',''), key=f"pdf_title_{di}")
                        doc['subtitle'] = st.text_input("Subtitle", value=doc.get('subtitle',''), key=f"pdf_subtitle_{di}")
                    with c2:
                        doc['type'] = st.text_input("Document Type", value=doc.get('type',''), key=f"pdf_type_{di}")
                        doc['classification'] = st.selectbox("Classification", 
                            ["Confidential","Internal","Public","Restricted"],
                            index=["Confidential","Internal","Public","Restricted"].index(doc.get('classification','Confidential'))
                            if doc.get('classification','Confidential') in ["Confidential","Internal","Public","Restricted"] else 0,
                            key=f"pdf_class_{di}")
                    
                    for si, sec in enumerate(doc.get('sections', [])):
                        st.markdown(f"**Section: {sec.get('title','')}**")
                        sec_cols = st.columns(2)
                        with sec_cols[0]:
                            if sec.get('intro') is not None:
                                sec['intro'] = st.text_area(f"Intro", value=sec.get('intro',''), 
                                    key=f"pdf_intro_{di}_{si}", height=80)
                            if sec.get('highlight') is not None:
                                sec['highlight'] = st.text_area(f"Highlight Quote", 
                                    value=sec.get('highlight',''), key=f"pdf_highlight_{di}_{si}", height=70)
                        with sec_cols[1]:
                            if sec.get('body') is not None:
                                sec['body'] = st.text_area(f"Body", value=sec.get('body',''),
                                    key=f"pdf_body_{di}_{si}", height=80)
                        
                        if sec.get('bullets'):
                            bullets_str = "\n".join(sec['bullets'])
                            new_b = st.text_area(f"Bullets (one per line)", value=bullets_str,
                                key=f"pdf_bullets_{di}_{si}", height=100)
                            sec['bullets'] = [b.strip() for b in new_b.split('\n') if b.strip()]
            
            st.session_state.pdf_docs = pdf_d
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 5 — SEO
    # ════════════════════════════════════════════════════════════════
    with tabs[4]:
        st.markdown('<div class="chip chip-teal">🔍 SEO METADATA</div>', unsafe_allow_html=True)
        
        seo = st.session_state.seo_data
        if not seo:
            st.info("🤖 Run SeoAgent to generate SEO metadata first.")
        else:
            col_s1, col_s2 = st.columns(2)
            with col_s1:
                seo['meta_title'] = st.text_input(f"Meta Title ({len(seo.get('meta_title',''))} chars)", 
                    value=seo.get('meta_title',''))
                seo['meta_description'] = st.text_area(f"Meta Description ({len(seo.get('meta_description',''))} chars)", 
                    value=seo.get('meta_description',''), height=80)
                seo['focus_keyword'] = st.text_input("Focus Keyword", value=seo.get('focus_keyword',''))
                seo['h1_recommendation'] = st.text_input("H1 SEO Variant", value=seo.get('h1_recommendation',''))
            with col_s2:
                seo['og_title'] = st.text_input("OG Title", value=seo.get('og_title',''))
                seo['og_description'] = st.text_area("OG Description", value=seo.get('og_description',''), height=80)
                seo['twitter_card'] = st.selectbox("Twitter Card", 
                    ["summary_large_image","summary"],
                    index=0 if seo.get('twitter_card') == 'summary_large_image' else 0)
            
            if seo.get('keywords'):
                st.markdown("**🏷️ Keywords**")
                keywords_str = ", ".join(seo['keywords'])
                kw_edit = st.text_input("Keywords (comma-separated)", value=keywords_str)
                seo['keywords'] = [k.strip() for k in kw_edit.split(',') if k.strip()]
            
            if seo.get('schema_json_ld'):
                with st.expander("📋 JSON-LD Schema Markup", expanded=False):
                    schema_str = json.dumps(seo['schema_json_ld'], indent=2)
                    st.code(schema_str, language="json")
                    st.markdown('<div class="tip-block">This JSON-LD is automatically injected into the HTML &lt;head&gt; when you export. It helps search engines understand your content.</div>', unsafe_allow_html=True)
            
            # SEO PREVIEW
            st.markdown("**🔍 SERP Preview**")
            meta_title = seo.get('meta_title', data.get('headline',''))
            meta_desc = seo.get('meta_description', data.get('subheadline',''))
            brand = data.get('brand_name', 'yourbrand')
            st.markdown(f"""
            <div style="background:#fff;border:1px solid #dfe1e5;border-radius:8px;padding:16px 20px;max-width:620px;font-family:'Arial',sans-serif">
                <div style="font-size:12px;color:#006621;margin-bottom:2px">{brand.lower().replace(' ','')}.com</div>
                <div style="font-size:18px;color:#1a0dab;margin-bottom:4px">{meta_title[:60]}</div>
                <div style="font-size:13px;color:#545454;line-height:1.5">{meta_desc[:160]}</div>
            </div>""", unsafe_allow_html=True)
            
            st.session_state.seo_data = seo
            # Also update page_data seo
            if st.session_state.page_data:
                st.session_state.page_data['seo_meta'] = seo
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 6 — QA REPORT
    # ════════════════════════════════════════════════════════════════
    with tabs[5]:
        st.markdown('<div class="chip chip-success">✅ QA REPORT</div>', unsafe_allow_html=True)
        
        qa = st.session_state.qa_data
        if not qa:
            st.info("🤖 Run QaAgent to get quality scoring and improvement suggestions first.")
        else:
            # Score display
            col_scores = st.columns(5)
            score_fields = [
                ("overall_score", "Overall", "🎯"),
                ("headline_score", "Headline", "📰"),
                ("cta_score", "CTA", "🔘"),
                ("copy_clarity_score", "Clarity", "📝"),
                ("conversion_score", "Conversion", "📈"),
            ]
            for i, (field, label, icon) in enumerate(score_fields):
                with col_scores[i]:
                    score = qa.get(field, 0)
                    color = "#2dd4a0" if score >= 80 else "#f59e0b" if score >= 60 else "#ff4466"
                    st.markdown(f"""
                    <div class="stat-box">
                        <div style="font-size:1.2rem;margin-bottom:4px">{icon}</div>
                        <div class="n" style="color:{color}">{score}</div>
                        <div class="l">{label}</div>
                    </div>""", unsafe_allow_html=True)
            
            st.markdown("<hr>", unsafe_allow_html=True)
            
            col_qa1, col_qa2 = st.columns(2)
            with col_qa1:
                if qa.get('strengths'):
                    st.markdown("**✅ Strengths**")
                    for s in qa['strengths']:
                        st.markdown(f"✓ {s}")
                
                if qa.get('improvements'):
                    st.markdown("**⚠️ Improvement Suggestions**")
                    for imp in qa['improvements']:
                        st.markdown(f"→ {imp}")
            
            with col_qa2:
                if qa.get('issues'):
                    st.markdown("**🔴 Issues Found**")
                    for issue in qa['issues']:
                        st.markdown(f"• {issue}")
                
                if qa.get('conversion_verdict'):
                    st.markdown("**🎯 Conversion Verdict**")
                    st.markdown(f'<div class="tip-block">{qa["conversion_verdict"]}</div>', unsafe_allow_html=True)
            
            if qa.get('ab_variants'):
                st.markdown("<hr>", unsafe_allow_html=True)
                st.markdown("**🧪 A/B Test Variants**")
                ab = qa['ab_variants']
                if ab.get('headline_variant_b'):
                    st.markdown(f"**Headline B:** {ab['headline_variant_b']}")
                if ab.get('cta_variant_b'):
                    st.markdown(f"**CTA B:** {ab['cta_variant_b']}")
                if ab.get('subheadline_variant_b'):
                    st.markdown(f"**Subheadline B:** {ab['subheadline_variant_b']}")
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 7 — EMAIL & ADS
    # ════════════════════════════════════════════════════════════════
    with tabs[6]:
        st.markdown('<div class="chip chip-success">📢 EMAIL & AD COPY</div>', unsafe_allow_html=True)
        
        copy_d = st.session_state.copy_data
        if not copy_d:
            st.info("🤖 Run CopyAgent to generate email sequences and ad copy first.")
        else:
            # Email subjects
            if copy_d.get('email_subject_lines'):
                with st.expander("📧 Email Subject Lines", expanded=True):
                    for sub in copy_d['email_subject_lines']:
                        st.markdown(f"• `{sub}`")
            
            # Email sequence
            if copy_d.get('email_sequence'):
                with st.expander("📬 Email Sequence (3 emails)", expanded=False):
                    for i, email in enumerate(copy_d['email_sequence']):
                        st.markdown(f"**Email {i+1} ({email.get('type','')})** — Subject: {email.get('subject','')}")
                        st.text_area(f"Email {i+1} Body", value=email.get('body',''), 
                                    key=f"email_body_{i}", height=120)
            
            # Ad copy
            if copy_d.get('ad_copy'):
                with st.expander("📣 Ad Copy (Google, Facebook, LinkedIn)", expanded=False):
                    for i, ad in enumerate(copy_d['ad_copy']):
                        st.markdown(f"**{ad.get('platform','')} — {ad.get('type','')}**")
                        c1, c2 = st.columns([1, 2])
                        with c1:
                            st.markdown(f"*Headline:* {ad.get('headline','')}")
                        with c2:
                            st.text_area(f"Ad Description", value=ad.get('description',''), 
                                        key=f"ad_desc_{i}", height=80)
                        st.markdown("<hr>", unsafe_allow_html=True)
            
            # Social posts
            if copy_d.get('social_posts'):
                with st.expander("📲 Social Media Posts", expanded=False):
                    for i, post in enumerate(copy_d['social_posts']):
                        st.markdown(f"**{post.get('platform','')}**")
                        st.text_area(f"Post Text", value=post.get('text',''), 
                                    key=f"post_text_{i}", height=100)
                        if post.get('hashtags'):
                            st.markdown(f"Hashtags: {' '.join(post['hashtags'])}")
                        st.markdown("<hr>", unsafe_allow_html=True)
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 8 — CODE
    # ════════════════════════════════════════════════════════════════
    with tabs[7]:
        st.markdown('<div class="chip chip-purple">💻 CODE OUTPUT</div>', unsafe_allow_html=True)
        
        html_out = generate_html_page(data, get_settings())
        lines = len(html_out.split('\n'))
        chars = len(html_out)
        
        st.markdown(f"""<div class="tip-block">
        📄 <strong>{lines} lines</strong> · {chars:,} characters · Standalone production HTML with embedded CSS, JS, SEO meta, and animations
        </div>""", unsafe_allow_html=True)
        
        col_c1, col_c2 = st.columns(2)
        with col_c1:
            st.markdown("**First 100 lines**")
            st.code('\n'.join(html_out.split('\n')[:100]), language="html")
        with col_c2:
            st.markdown("**Page structure**")
            sections = ["<!DOCTYPE html>", "<head>", "<style>", "</style>", "</head>", 
                       "<body>", "<nav>", ".hero", ".stats", ".feat", ".test", ".faq", 
                       ".cband", "<footer>", "<script>", "</html>"]
            for s in sections:
                count = html_out.count(s)
                if count:
                    st.markdown(f"`{s}` — {count}×")
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 9 — PREVIEW
    # ════════════════════════════════════════════════════════════════
    with tabs[8]:
        st.markdown('<div class="chip chip-red">👁️ LIVE PREVIEW</div>', unsafe_allow_html=True)
        
        preview_settings = get_settings()
        preview_settings['nav_sticky'] = False
        preview_html = generate_html_page(data, preview_settings)
        
        # Device selector
        dcol1, dcol2 = st.columns([1, 3])
        with dcol1:
            device = st.selectbox("Preview Device", ["Desktop (1160px)", "Tablet (768px)", "Mobile (375px)"])
        
        device_widths = {
            "Desktop (1160px)": "100%",
            "Tablet (768px)": "768px",
            "Mobile (375px)": "375px"
        }
        preview_w = device_widths.get(device, "100%")
        
        url_brand = data.get('brand_name','yourbrand').lower().replace(' ','')
        st.markdown(f"""
        <div class="preview-browser-bar">
            <div class="browser-dot" style="background:#ff5f57"></div>
            <div class="browser-dot" style="background:#febc2e"></div>
            <div class="browser-dot" style="background:#28c840"></div>
            <div class="browser-url">{url_brand}.com &nbsp;·&nbsp; {device}</div>
        </div>""", unsafe_allow_html=True)
        
        if preview_w == "100%":
            st.components.v1.html(preview_html, height=960, scrolling=True)
        else:
            st.markdown(f'<div style="display:flex;justify-content:center;border:1px solid var(--border,#1a1e32);border-top:none;border-radius:0 0 8px 8px;background:#fff">', unsafe_allow_html=True)
            st.components.v1.html(
                f'<div style="width:{preview_w};margin:0 auto;overflow-x:hidden">{preview_html}</div>',
                height=960, scrolling=True
            )
            st.markdown('</div>', unsafe_allow_html=True)
    
    
    # ════════════════════════════════════════════════════════════════
    # TAB 10 — EXPORT
    # ════════════════════════════════════════════════════════════════
    with tabs[9]:
        st.markdown('<div class="chip chip-success">📦 EXPORT & DOWNLOAD</div>', unsafe_allow_html=True)
        
        settings_final = get_settings()
        html_final = generate_html_page(data, settings_final)
        ts = datetime.now().strftime('%Y%m%d-%H%M%S')
        brand_slug = re.sub(r'[^a-z0-9]', '-', data.get('brand_name','brand').lower())
        
        # Summary row
        st.markdown(f"""
        <div class="docs-summary-row">
            <div class="doc-badge"><div class="doc-badge-icon">📄</div><div class="doc-badge-count">1</div><div class="doc-badge-label">HTML Page</div></div>
            <div class="doc-badge"><div class="doc-badge-icon">📊</div><div class="doc-badge-count">{len(st.session_state.pptx_data.get('slides',[]) if st.session_state.pptx_data else [])}</div><div class="doc-badge-label">PPTX Slides</div></div>
            <div class="doc-badge"><div class="doc-badge-icon">📁</div><div class="doc-badge-count">{len(st.session_state.pdf_docs) if st.session_state.pdf_docs else 0}</div><div class="doc-badge-label">PDF Docs</div></div>
            <div class="doc-badge"><div class="doc-badge-icon">🔍</div><div class="doc-badge-count">{'✓' if st.session_state.seo_data else '—'}</div><div class="doc-badge-label">SEO Meta</div></div>
            <div class="doc-badge"><div class="doc-badge-icon">📢</div><div class="doc-badge-count">{'✓' if st.session_state.copy_data else '—'}</div><div class="doc-badge-label">Ad Copy</div></div>
        </div>""", unsafe_allow_html=True)
        
        st.markdown("<hr>", unsafe_allow_html=True)
        
        col_e1, col_e2, col_e3, col_e4 = st.columns(4)
        
        # ── HTML DOWNLOAD ──
        with col_e1:
            st.markdown(f"""<div class="export-card">
                <div class="export-card-icon">📄</div>
                <div class="export-card-title">Landing Page HTML</div>
                <div class="export-card-meta">{len(html_final.split(chr(10)))} lines · {len(html_final):,} chars</div>
                <div class="export-card-desc">Standalone HTML file with embedded CSS, animations, SEO tags, and JavaScript. Upload to any static host.</div>
            </div>""", unsafe_allow_html=True)
            st.download_button(
                "⬇ Download HTML",
                data=html_final,
                file_name=f"{brand_slug}-landing-page-{ts}.html",
                mime="text/html",
                use_container_width=True,
                key="dl_html"
            )
        
        # ── PPTX DOWNLOAD ──
        with col_e2:
            pptx_slides = st.session_state.pptx_data.get('slides', []) if st.session_state.pptx_data else []
            pptx_font = st.session_state.get("pptx_font", "Calibri")
            
            if pptx_slides:
                try:
                    pptx_bytes = generate_pptx(
                        pptx_slides,
                        brand_hex=st.session_state.brand_color,
                        accent_hex=st.session_state.accent_color,
                        company_name=data.get('brand_name','Brand'),
                        font_header=pptx_font
                    )
                    st.markdown(f"""<div class="export-card">
                        <div class="export-card-icon">📊</div>
                        <div class="export-card-title">PowerPoint Deck</div>
                        <div class="export-card-meta">{len(pptx_slides)} slides · python-pptx</div>
                        <div class="export-card-desc">Branded .pptx file with custom color scheme, {pptx_font} typography, speaker notes, and professional layouts.</div>
                    </div>""", unsafe_allow_html=True)
                    st.download_button(
                        "⬇ Download .pptx",
                        data=pptx_bytes,
                        file_name=f"{brand_slug}-deck-{ts}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                        key="dl_pptx"
                    )
                except Exception as e:
                    st.error(f"PPTX error: {e}")
            else:
                st.markdown("""<div class="export-card">
                    <div class="export-card-icon">📊</div>
                    <div class="export-card-title">PowerPoint Deck</div>
                    <div class="export-card-desc">Run PptxAgent to generate slide content first.</div>
                </div>""", unsafe_allow_html=True)
        
        # ── PDF ZIP DOWNLOAD ──
        with col_e3:
            pdf_docs = st.session_state.pdf_docs
            if pdf_docs:
                try:
                    pdf_zip = generate_all_pdfs(
                        pdf_docs,
                        brand_hex=st.session_state.brand_color,
                        company_name=data.get('brand_name','Brand')
                    )
                    st.markdown(f"""<div class="export-card">
                        <div class="export-card-icon">📁</div>
                        <div class="export-card-title">PDF Document Bundle</div>
                        <div class="export-card-meta">{len(pdf_docs)} PDFs · ZIP archive · reportlab</div>
                        <div class="export-card-desc">All {len(pdf_docs)} documents in a ZIP: Executive Brief, Sales Proposal, Product Overview, and One-Pager.</div>
                    </div>""", unsafe_allow_html=True)
                    st.download_button(
                        f"⬇ Download {len(pdf_docs)} PDFs (ZIP)",
                        data=pdf_zip,
                        file_name=f"{brand_slug}-docs-{ts}.zip",
                        mime="application/zip",
                        use_container_width=True,
                        key="dl_pdf_zip"
                    )
                    
                    # Individual PDFs
                    with st.expander("⬇ Download Individual PDFs", expanded=False):
                        for di, doc in enumerate(pdf_docs):
                            try:
                                pdf_b = generate_pdf_document(
                                    doc, 
                                    brand_hex=st.session_state.brand_color,
                                    company_name=data.get('brand_name','Brand')
                                )
                                safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', doc.get('type', f'doc_{di}'))
                                st.download_button(
                                    f"⬇ {doc.get('type','Document')}",
                                    data=pdf_b,
                                    file_name=f"{brand_slug}-{safe_name}-{ts}.pdf",
                                    mime="application/pdf",
                                    use_container_width=True,
                                    key=f"dl_pdf_{di}"
                                )
                            except Exception as e:
                                st.error(f"PDF {di} error: {e}")
                except Exception as e:
                    st.error(f"PDF generation error: {e}")
            else:
                st.markdown("""<div class="export-card">
                    <div class="export-card-icon">📁</div>
                    <div class="export-card-title">PDF Document Bundle</div>
                    <div class="export-card-desc">Run PdfAgent to generate the document suite first.</div>
                </div>""", unsafe_allow_html=True)
        
        # ── FULL BUNDLE / JSON ──
        with col_e4:
            # Build full content bundle JSON
            bundle = {
                "page_content": st.session_state.page_data,
                "pptx_slides": st.session_state.pptx_data,
                "pdf_docs": st.session_state.pdf_docs,
                "seo_data": st.session_state.seo_data,
                "qa_report": st.session_state.qa_data,
                "copy_data": st.session_state.copy_data,
                "design_settings": {
                    "brand_color": st.session_state.brand_color,
                    "accent_color": st.session_state.accent_color,
                    "bg_color": st.session_state.bg_color,
                    "text_color": st.session_state.text_color,
                    "font_pair": st.session_state.font_pair,
                    "hero_layout": st.session_state.hero_layout,
                },
                "generated_at": datetime.now().isoformat(),
                "total_tokens": orch.total_tokens,
            }
            bundle_json = json.dumps(bundle, indent=2)
            
            st.markdown(f"""<div class="export-card">
                <div class="export-card-icon">🗃️</div>
                <div class="export-card-title">Full Content Bundle</div>
                <div class="export-card-meta">{len(bundle_json):,} chars · JSON</div>
                <div class="export-card-desc">All generated content — page copy, slides, PDF docs, SEO, ad copy, QA report — as structured JSON.</div>
            </div>""", unsafe_allow_html=True)
            st.download_button(
                "⬇ Download Bundle JSON",
                data=bundle_json,
                file_name=f"{brand_slug}-full-bundle-{ts}.json",
                mime="application/json",
                use_container_width=True,
                key="dl_bundle"
            )
            
            st.markdown("<hr>", unsafe_allow_html=True)
            if st.button("🔄 Regenerate Everything", use_container_width=True):
                if api_key and st.session_state.description:
                    st.session_state.page_data = None
                    st.session_state.pptx_data = None
                    st.session_state.pdf_docs = None
                    st.session_state.seo_data = None
                    st.session_state.design_data = None
                    st.session_state.qa_data = None
                    st.session_state.copy_data = None
                    st.rerun()
        
        # ── AGENT LOG ──
        st.markdown("<hr>", unsafe_allow_html=True)
        st.markdown('<div class="chip chip-purple">🤖 AGENT EXECUTION LOG</div>', unsafe_allow_html=True)
        
        if st.session_state.agent_log:
            log_html = '<div style="background:#0c0e1a;border:1px solid #1a1e32;border-radius:8px;padding:12px;max-height:300px;overflow-y:auto;font-family:\'Space Mono\',monospace;font-size:0.68rem;">'
            for entry in st.session_state.agent_log:
                s = entry.get("status","info")
                c = {"running":"#f59e0b","done":"#2dd4a0","error":"#ff4466","info":"#4d6fff"}.get(s,"#525870")
                log_html += f'<div style="padding:4px 0;border-bottom:1px solid rgba(255,255,255,0.04);"><span style="color:#2a2f48">{entry.get("ts","")} </span><span style="color:{c};font-weight:700">[{entry.get("agent","")}]</span> <span style="color:#e4e7f0">{entry.get("msg","")}</span></div>'
            log_html += f'<div style="margin-top:8px;color:#2a2f48">Total tokens: {orch.total_tokens} · Agents run: {sum(1 for a in orch.agents.values() if a.status=="done")}/{len(orch.agents)}</div>'
            log_html += "</div>"
            st.markdown(log_html, unsafe_allow_html=True)
        else:
            st.markdown('<span style="color:#3a3f58;font-family:\'Space Mono\',monospace;font-size:0.72rem">No agent runs yet — launch the swarm to see logs here</span>', unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────
# EMPTY STATE
# ─────────────────────────────────────────────────────────────────

else:
    st.markdown("""
    <div style="text-align:center;padding:4rem 2rem 2rem">
        <div style="font-size:3.5rem;margin-bottom:1rem">⚡</div>
        <div style="font-family:'Bebas Neue',sans-serif;font-size:2.2rem;letter-spacing:0.08em;color:#e4e7f0;margin-bottom:8px">
            READY TO FORGE
        </div>
        <p style="font-family:'Space Mono',monospace;font-size:0.7rem;color:#3a3f58;letter-spacing:0.14em;text-transform:uppercase;max-width:480px;margin:0 auto 2rem">
            Enter your business description · Configure your agents · Launch the swarm
        </p>
    </div>
    <div style="display:grid;grid-template-columns:repeat(4,1fr);gap:12px;max-width:900px;margin:0 auto">
    """, unsafe_allow_html=True)
    
    outputs = [
        ("📄", "Landing Page HTML", "Full standalone page with hero, features, testimonials, FAQ, animations & SEO"),
        ("📊", "PowerPoint Deck", "Real .pptx with custom brand colors, fonts, layouts and speaker notes"),
        ("📁", "PDF Document Suite", "Executive brief, sales proposal, product overview & one-pager as PDFs"),
        ("🔍", "SEO Metadata", "Meta tags, OpenGraph, JSON-LD schema & keyword strategy"),
        ("🎨", "Design System", "AI-recommended colors, typography, layout & animation style"),
        ("✅", "QA Report", "Quality scoring with A/B variants and conversion optimization tips"),
        ("📢", "Email & Ad Copy", "Email sequences, Google/Facebook/LinkedIn ad copy & social posts"),
        ("🗃️", "Content Bundle", "Complete JSON export of all generated content for reuse"),
    ]
    
    oc = st.columns(4)
    for i, (icon, title, desc) in enumerate(outputs):
        with oc[i % 4]:
            st.markdown(f"""<div class="export-card" style="margin-bottom:12px">
                <div class="export-card-icon">{icon}</div>
                <div class="export-card-title">{title}</div>
                <div class="export-card-desc">{desc}</div>
            </div>""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────
# FOOTER
# ─────────────────────────────────────────────────────────────────

st.markdown("<hr>", unsafe_allow_html=True)
st.markdown(f"""
<div style="display:flex;justify-content:space-between;align-items:center;padding:0.5rem 0;flex-wrap:wrap;gap:8px">
    <span style="font-family:'Space Mono',monospace;font-size:0.6rem;color:#1e2335;letter-spacing:0.16em;text-transform:uppercase">
        PAGEFORGE PRO — AI DOCUMENT STUDIO — BMADD AGENT FRAMEWORK v2.0
    </span>
    <span style="font-family:'Space Mono',monospace;font-size:0.6rem;color:#1e2335;letter-spacing:0.1em;text-transform:uppercase">
        HTML · PPTX · PDF · SEO · COPYWRITING — POWERED BY OPENAI + STREAMLIT + PYTHON-PPTX + REPORTLAB
    </span>
</div>
""", unsafe_allow_html=True)
